"""Builds the executive deck for the AHDS payer-isolation POC.

No customer or individual names appear anywhere in the output by design.

    py -3.14-arm64 -m venv %LOCALAPPDATA%\\venvs\\pptxbuild
    %LOCALAPPDATA%\\venvs\\pptxbuild\\Scripts\\python.exe build-deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

W, H = 13.333, 7.5

NAVY = RGBColor(0x0E, 0x2A, 0x47)
NAVY_CARD = RGBColor(0x1A, 0x3E, 0x60)
TEAL = RGBColor(0x1C, 0x72, 0x93)
TEAL_LT = RGBColor(0x4F, 0xA6, 0xC4)
ICE = RGBColor(0xCF, 0xE3, 0xEF)
PAPER = RGBColor(0xF5, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x12, 0x28, 0x3D)
GREY = RGBColor(0x5A, 0x6B, 0x7B)
LINE = RGBColor(0xD8, 0xE2, 0xEA)
AMBER = RGBColor(0xD9, 0x82, 0x2B)
AMBER_LT = RGBColor(0xFD, 0xF3, 0xE3)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
GREEN_LT = RGBColor(0xE8, 0xF4, 0xEC)
RED = RGBColor(0xB0, 0x3A, 0x2B)
RED_LT = RGBColor(0xFA, 0xEC, 0xEA)

HEAD = "Cambria"
BODY = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

_n = {"i": 0}


def new(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY if dark else PAPER
    _n["i"] += 1
    return s


def rect(s, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE, radius=None,
         line=None, lw=1.0):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def txt(s, x, y, w, h, runs, size=14, color=INK, font=BODY, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, italic=False,
        line_spacing=1.0):
    """runs: str, or list of (text, {overrides}) tuples, or list of str lines."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    if isinstance(runs, str):
        runs = [(runs, {})]
    runs = [(r, {}) if isinstance(r, str) else r for r in runs]

    for i, (text, ov) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("line_spacing", line_spacing)
        if i > 0:
            p.space_before = Pt(ov.get("space_before", space))
        if ov.get("bullet_gap"):
            p.space_before = Pt(ov["bullet_gap"])
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic)
        f.color.rgb = ov.get("color", color)
    return tb


def head(s, kicker, title, dark=False):
    """Standard content-slide header: accent bar, kicker, title, footer rule."""
    rect(s, 0.62, 0.46, 0.075, 0.72, fill=TEAL_LT if dark else TEAL)
    txt(s, 0.85, 0.46, 11.8, 0.26, kicker.upper(), size=10.5, bold=True,
        color=TEAL_LT if dark else TEAL, font=BODY)
    txt(s, 0.85, 0.72, 11.8, 0.52, title, size=27, bold=True,
        color=WHITE if dark else INK, font=HEAD)
    foot(s, dark)


def foot(s, dark=False):
    rect(s, 0.62, 6.94, 12.1, 0.014,
         fill=RGBColor(0x2C, 0x50, 0x74) if dark else LINE)
    txt(s, 0.62, 7.03, 9.0, 0.24,
        "Azure Health Data Services  ·  Payer data exchange reference implementation",
        size=9, color=RGBColor(0x7F, 0x9A, 0xB5) if dark else GREY)
    txt(s, 11.0, 7.03, 1.72, 0.24, f"{_n['i']:02d}", size=9, align=PP_ALIGN.RIGHT,
        color=RGBColor(0x7F, 0x9A, 0xB5) if dark else GREY)


def card(s, x, y, w, h, fill=WHITE, line=LINE, radius=0.06):
    return rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                radius=radius, line=line, lw=0.75)


def chip(s, x, y, d, label, fill=TEAL, color=WHITE, size=12):
    rect(s, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)
    txt(s, x, y + d / 2 - 0.115, d, 0.25, label, size=size, bold=True,
        color=color, align=PP_ALIGN.CENTER, font=BODY)


def stat(s, x, y, w, value, label, color=TEAL, vsize=34, dark=False):
    txt(s, x, y, w, 0.62, value, size=vsize, bold=True, color=color, font=HEAD)
    txt(s, x, y + 0.62, w, 0.5, label, size=11,
        color=RGBColor(0x9F, 0xB8, 0xCC) if dark else GREY)


# ---------------------------------------------------------------- 01  title
s = new(dark=True)
rect(s, 0, 0, 0.28, H, fill=TEAL)
rect(s, 7.9, 0, 5.44, H, fill=RGBColor(0x0B, 0x22, 0x3A))
for i, o in enumerate([0.0, 0.62, 1.24]):
    rect(s, 8.55 + o, 1.62, 0.44, 3.6, fill=RGBColor(0x14, 0x33, 0x54))
    rect(s, 8.55 + o, 1.62, 0.44, 0.44, fill=TEAL if i == 0 else NAVY_CARD)
rect(s, 10.9, 1.62, 1.9, 3.6, fill=RGBColor(0x14, 0x33, 0x54))
rect(s, 10.9, 1.62, 1.9, 0.44, fill=TEAL_LT)
txt(s, 8.55, 5.36, 4.3, 0.9,
    [("One gateway.  Two isolated data planes.", {"bold": True, "color": ICE}),
     ("No payer credential ever reaches the FHIR service.",
      {"color": RGBColor(0x7F, 0x9A, 0xB5), "space_before": 4})],
    size=10.5)

txt(s, 1.0, 1.9, 6.6, 0.3, "AZURE HEALTH DATA SERVICES", size=11.5, bold=True,
    color=TEAL_LT)
txt(s, 1.0, 2.32, 6.6, 1.9,
    [("Payer data isolation,", {}), ("proven — not asserted.", {"color": ICE})],
    size=42, bold=True, color=WHITE, font=HEAD, line_spacing=1.02)
rect(s, 1.0, 4.34, 1.5, 0.035, fill=TEAL)
txt(s, 1.0, 4.68, 6.4, 1.0,
    "A deployed reference implementation for CMS-0057-F payer data exchange: "
    "isolation at the payer, contract scoping at the gateway, and an executable "
    "test suite that fails loudly when either stops being true.",
    size=13.5, color=RGBColor(0xB9, 0xCD, 0xDD), line_spacing=1.28)
txt(s, 1.0, 6.1, 6.4, 0.3,
    "16 of 16 isolation assertions passing  ·  live environment  ·  16 August 2026",
    size=10.5, bold=True, color=TEAL_LT)

# ---------------------------------------------------------------- 02  problem
s = new()
head(s, "Why this is hard", "Forty payers. One boundary that cannot be wrong.")
for i, (v, l) in enumerate([("~40", "payer partners at target scale"),
                            ("~800K", "patients in scope"),
                            ("4", "required API families"),
                            ("Jan 2027", "CMS-0057-F enforcement")]):
    x = 0.62 + i * 3.06
    card(s, x, 1.62, 2.86, 1.68)
    stat(s, x + 0.3, 1.85, 2.4, v, l, vsize=30 if v != "Jan 2027" else 24)

card(s, 0.62, 3.62, 5.94, 2.98, fill=WHITE)
rect(s, 0.62, 3.62, 0.06, 2.98, fill=RED)
txt(s, 1.0, 3.94, 5.2, 0.3, "THE FAILURE THAT MATTERS", size=10, bold=True, color=RED)
txt(s, 1.0, 4.34, 5.2, 2.1,
    "One payer reaching another payer's PHI is not a bug report. It is a "
    "reportable disclosure, a contractual breach, and a regulatory event — all "
    "at once, and none of them reversible.\n\nEverything else on this list is "
    "recoverable. This one is not.",
    size=13, color=INK, line_spacing=1.32)

card(s, 6.78, 3.62, 5.94, 2.98)
rect(s, 6.78, 3.62, 0.06, 2.98, fill=TEAL)
txt(s, 7.16, 3.94, 5.2, 0.3, "SO THE BOUNDARY MUST SURVIVE", size=10, bold=True,
    color=TEAL)
for i, t in enumerate(["A defect in a gateway policy",
                       "A well-intentioned role assignment made while debugging",
                       "A payer who has read the documentation and found the direct endpoint",
                       "Four years and a change of staff"]):
    chip(s, 7.16, 4.42 + i * 0.55, 0.16, "", fill=TEAL)
    txt(s, 7.5, 4.36 + i * 0.55, 4.96, 0.5, t, size=12, color=INK)

# ---------------------------------------------------------------- 03  the ask
s = new()
head(s, "The ask", "Nine questions from the working session")
asks = [
    ("01", "Reference architecture", "A diagram we can put in front of security review."),
    ("02", "Where does isolation go?", "Physical per payer, or logical filtering?"),
    ("03", "PHI at the payer level", "Sharing agreements are per payer. Separation must match."),
    ("04", "Keep payers out of ingest", "A payer must not query our inbound data."),
    ("05", "Is a gateway mandatory?", "Or can the FHIR service enforce this alone?"),
    ("06", "How do we onboard a payer?", "The exact process, start to finish."),
    ("07", "Why is $import failing?", "A 403 blocking the development environment."),
    ("08", "Quota: 10 to 40 services", "What does raising it involve?"),
    ("09", "Concurrent export capacity", "Forty payers pull at once. Does it hold?"),
]
for i, (n, t, d) in enumerate(asks):
    c, r = i % 3, i // 3
    x, y = 0.62 + c * 4.08, 1.66 + r * 1.78
    card(s, x, y, 3.88, 1.6)
    txt(s, x + 0.3, y + 0.26, 0.6, 0.3, n, size=15, bold=True, color=TEAL_LT, font=HEAD)
    txt(s, x + 0.3, y + 0.62, 3.3, 0.32, t, size=13.5, bold=True, color=INK)
    txt(s, x + 0.3, y + 0.98, 3.3, 0.5, d, size=10.5, color=GREY, line_spacing=1.2)

# ---------------------------------------------------------------- 04  response
s = new(dark=True)
head(s, "The response", "Not a deck. A deployed environment with an executable proof.",
     dark=True)
for i, (v, l) in enumerate([("2", "isolated FHIR services, one per payer"),
                            ("16 / 16", "isolation assertions passing"),
                            ("7", "onboarding steps, fully scripted"),
                            ("~20 min", "to rebuild the whole thing from zero")]):
    x = 0.62 + i * 3.06
    rect(s, x, 1.72, 2.86, 1.76, fill=NAVY_CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.06)
    rect(s, x, 1.72, 2.86, 0.05, fill=TEAL_LT)
    stat(s, x + 0.32, 2.0, 2.3, v, l, color=WHITE,
         vsize=30 if len(v) < 6 else 24, dark=True)

txt(s, 0.62, 3.94, 12.1, 0.4,
    "Every claim in this deck maps to a command that runs against the "
    "subscription today.", size=17, bold=True, color=ICE, font=HEAD)

items = [
    ("Infrastructure as code", "One Bicep template. Adding the fortieth payer is one array entry."),
    ("Gateway policy", "Six enforcement layers, deny by default, version-controlled."),
    ("Proof suite", "Sixteen assertions against the live gateway. Three minutes."),
    ("Documentation", "Eight technical notes, two runbooks, an editable architecture set."),
]
for i, (t, d) in enumerate(items):
    x, y = 0.62 + (i % 2) * 6.16, 4.6 + (i // 2) * 1.06
    rect(s, x, y + 0.08, 0.05, 0.7, fill=TEAL)
    txt(s, x + 0.3, y + 0.06, 5.5, 0.3, t, size=13, bold=True, color=WHITE)
    txt(s, x + 0.3, y + 0.4, 5.5, 0.5, d, size=11,
        color=RGBColor(0x9F, 0xB8, 0xCC), line_spacing=1.2)

# ---------------------------------------------------------------- 05  options
s = new()
head(s, "Decision", "Where the hard boundary goes")
opts = [
    ("Option 1", "One service per payer.\nContracts separated logically.",
     "CHOSEN", GREEN, GREEN_LT,
     "PHI boundary matches the sharing agreement. ~40 instances, template-managed."),
    ("Option 2", "One shared service.\nEverything separated logically.",
     "REJECTED", RED, RED_LT,
     "Co-mingles every payer's PHI in one store. A single policy defect is a cross-payer disclosure."),
    ("Option 3", "One service per contract.\nPhysical all the way down.",
     "REJECTED", RED, RED_LT,
     "Doubles the instance count to ~80+ and doubles the quota problem. Buys containment nobody needs."),
    ("Option 4", "Shared inbound service,\nsplit outbound services.",
     "REJECTED", RED, RED_LT,
     "The shared inbound store still co-mingles PHI — the exact condition ruled out."),
]
for i, (n, d, badge, bc, bl, why) in enumerate(opts):
    x = 0.62 + i * 3.06
    card(s, x, 1.62, 2.86, 3.5, fill=WHITE, line=GREEN if i == 0 else LINE)
    rect(s, x, 1.62, 2.86, 0.055, fill=GREEN if i == 0 else RGBColor(0xC6, 0xD3, 0xDC))
    txt(s, x + 0.28, 1.9, 2.3, 0.28, n, size=12, bold=True, color=GREY)
    txt(s, x + 0.28, 2.24, 2.34, 0.9, d, size=13.5, bold=True, color=INK,
        line_spacing=1.22)
    rect(s, x + 0.28, 3.24, 1.16, 0.3, fill=bl, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.3)
    txt(s, x + 0.28, 3.31, 1.16, 0.24, badge, size=9, bold=True, color=bc,
        align=PP_ALIGN.CENTER)
    txt(s, x + 0.28, 3.72, 2.34, 1.2, why, size=10.5, color=GREY, line_spacing=1.24)

rect(s, 0.62, 5.42, 12.1, 1.16, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.05)
txt(s, 1.05, 5.66, 11.3, 0.34, "Physical separation at the payer.  Logical separation at the contract.",
    size=19, bold=True, color=WHITE, font=HEAD)
txt(s, 1.05, 6.08, 11.3, 0.3,
    "This is the split that makes roughly forty instances defensible instead of two hundred.",
    size=11.5, color=RGBColor(0x9F, 0xB8, 0xCC))

# ---------------------------------------------------------------- 06  two boundaries
s = new()
head(s, "The reasoning", "Two boundaries doing two different jobs")
cols = [
    ("PHYSICAL — per payer", TEAL,
     "Limits the blast radius of a mistake.",
     [("Enforced by", "Azure RBAC. There is no authorisation path from one payer to another payer's service."),
      ("Failure mode", "Cross-payer PHI disclosure."),
      ("Recoverable?", "No. Contractual and regulatory, immediately."),
      ("Cost of getting it wrong", "Unbounded.")]),
    ("LOGICAL — per contract", RGBColor(0x6B, 0x4C, 0x9A),
     "Limits the result set of a query.",
     [("Enforced by", "meta.tag stamped on write; a forced _tag filter injected on every read."),
      ("Failure mode", "Over-disclosure inside one payer's own data."),
      ("Recoverable?", "Yes. Contained, auditable, correctable."),
      ("Cost of getting it wrong", "Embarrassing, not existential.")]),
]
for i, (title, col, sub, rows) in enumerate(cols):
    x = 0.62 + i * 6.16
    card(s, x, 1.62, 5.94, 4.34)
    rect(s, x, 1.62, 5.94, 0.055, fill=col)
    txt(s, x + 0.34, 1.9, 5.2, 0.28, title, size=11.5, bold=True, color=col)
    txt(s, x + 0.34, 2.24, 5.2, 0.34, sub, size=16, bold=True, color=INK, font=HEAD)
    for j, (k, v) in enumerate(rows):
        y = 2.78 + j * 0.79
        txt(s, x + 0.34, y, 1.72, 0.5, k, size=10, bold=True, color=GREY)
        txt(s, x + 2.16, y, 3.4, 0.72, v, size=11.5, color=INK, line_spacing=1.2)
        if j < 3:
            rect(s, x + 0.34, y + 0.66, 5.24, 0.01, fill=LINE)

rect(s, 0.62, 6.16, 12.1, 0.62, fill=ICE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.08)
txt(s, 1.0, 6.34, 11.4, 0.3,
    "Put the hard boundary where a failure is unrecoverable. Accept a soft boundary where a failure is contained.",
    size=13, bold=True, color=NAVY, font=HEAD)

# ---------------------------------------------------------------- 07  architecture
s = new()
head(s, "Reference architecture", "One way in. One way out. Nothing else routes.")

rect(s, 0.62, 1.78, 2.5, 1.62, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.07, line=LINE, lw=0.75)
txt(s, 0.86, 2.02, 2.05, 0.28, "PAYER PARTNERS", size=9.5, bold=True, color=TEAL)
txt(s, 0.86, 2.36, 2.05, 0.9,
    "Client credentials\nfrom Microsoft Entra ID\nScoped per contract",
    size=11, color=INK, line_spacing=1.24)

rect(s, 0.62, 3.72, 2.5, 1.42, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.07, line=LINE, lw=0.75)
txt(s, 0.86, 3.94, 2.05, 0.28, "INGEST PIPELINE", size=9.5, bold=True,
    color=RGBColor(0x6B, 0x4C, 0x9A))
txt(s, 0.86, 4.28, 2.05, 0.7, "Provider-side writes\nStamped with contract tag",
    size=11, color=INK, line_spacing=1.24)

rect(s, 4.06, 1.78, 2.92, 4.4, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.05)
txt(s, 4.36, 2.0, 2.3, 0.3, "API MANAGEMENT", size=10, bold=True, color=TEAL_LT)
txt(s, 4.36, 2.32, 2.34, 0.34, "The only ingress", size=15, bold=True, color=WHITE,
    font=HEAD)
for j, lab in enumerate(["Token validation", "Entitlement lookup",
                         "Route allow-list", "Contract scope injection",
                         "Rate + concurrency", "Managed-identity swap"]):
    y = 2.86 + j * 0.5
    rect(s, 4.36, y, 2.3, 0.4, fill=NAVY_CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.12)
    txt(s, 4.5, y + 0.1, 2.1, 0.24, lab, size=9.5, color=ICE)
txt(s, 4.36, 5.88, 2.34, 0.24, "Payer token stops here.", size=9.5, bold=True,
    italic=True, color=TEAL_LT)

for j, (nm, sub, y) in enumerate([("FHIR service — payer A", "contracts CT-3456 · CT-7788", 1.78),
                                  ("FHIR service — payer B", "contract CT-9001", 3.14)]):
    card(s, 7.92, y, 4.8, 1.16)
    rect(s, 7.92, y, 0.055, 1.16, fill=TEAL)
    txt(s, 8.26, y + 0.24, 4.2, 0.3, nm, size=13, bold=True, color=INK)
    txt(s, 8.26, y + 0.58, 4.2, 0.3, sub, size=10.5, color=GREY, font=MONO)

card(s, 7.92, 4.5, 2.32, 1.68)
txt(s, 8.22, 4.74, 1.9, 0.28, "STORAGE", size=9.5, bold=True, color=TEAL)
txt(s, 8.22, 5.06, 1.9, 0.9, "Export output\nImport staging\nQuarantine",
    size=10.5, color=INK, line_spacing=1.26)
card(s, 10.4, 4.5, 2.32, 1.68)
txt(s, 10.7, 4.74, 1.9, 0.28, "LOG ANALYTICS", size=9.5, bold=True, color=TEAL)
txt(s, 10.7, 5.06, 1.9, 0.9, "Gateway denials\nStorage 403s\nExport durations",
    size=10.5, color=INK, line_spacing=1.26)

for y0, y1 in [(2.5, 2.5), (4.34, 4.34)]:
    rect(s, 3.12, y0, 0.94, 0.02, fill=TEAL)
rect(s, 6.98, 2.36, 0.94, 0.02, fill=TEAL)
rect(s, 6.98, 3.72, 0.94, 0.02, fill=TEAL)
txt(s, 3.12, 2.16, 0.94, 0.24, "read", size=9, color=GREY, align=PP_ALIGN.CENTER)
txt(s, 3.12, 4.0, 0.94, 0.24, "write", size=9, color=GREY, align=PP_ALIGN.CENTER)

txt(s, 0.62, 6.36, 12.1, 0.34,
    "Nothing reaches a FHIR service except through the gateway — because nothing else has been given a role that lets it.",
    size=12.5, bold=True, color=NAVY, font=HEAD)

# ---------------------------------------------------------------- 08  control plane
s = new()
head(s, "Control plane", "Six layers, in this order, deny by default")
layers = [
    ("1", "Validate the token", "Issuer, audience and expiry. A token minted for a different service is refused before anything else runs."),
    ("2", "Look up the entitlement", "Which payer is this, and which contracts do they hold? Unknown caller, unknown answer — denied."),
    ("3", "Apply the route allow-list", "Six operations are named. Everything not named is denied, including operations that do not exist yet."),
    ("4", "Inject the scope", "Group membership checked on export. A _tag filter forced onto every search, overwriting anything the caller supplied."),
    ("5", "Limit rate and concurrency", "600 requests per minute, 50,000 per day, and exactly one bulk export per payer per five minutes."),
    ("6", "Swap to the gateway identity", "The payer's token is discarded and replaced with the gateway's own managed identity token."),
]
for i, (n, t, d) in enumerate(layers):
    x, y = 0.62 + (i % 2) * 6.16, 1.66 + (i // 2) * 1.42
    card(s, x, y, 5.94, 1.26)
    chip(s, x + 0.3, y + 0.34, 0.44, n, fill=TEAL, size=14)
    txt(s, x + 0.96, y + 0.24, 4.72, 0.3, t, size=13.5, bold=True, color=INK)
    txt(s, x + 0.96, y + 0.58, 4.72, 0.62, d, size=10.5, color=GREY,
        line_spacing=1.24)

rect(s, 0.62, 6.06, 12.1, 0.64, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.07)
txt(s, 1.0, 6.24, 11.4, 0.3,
    "The payer's credential never reaches the FHIR service. The gateway is a broker, not a proxy.",
    size=13, bold=True, color=WHITE, font=HEAD)

# ---------------------------------------------------------------- 09  direction
s = new()
head(s, "Direction", "A payer cannot see what you ingest")
rows = [
    ("Who calls it", "The provider ingest pipeline", "The payer"),
    ("Verbs allowed", "GET · POST · PUT · DELETE", "GET only"),
    ("Bulk export", "Denied outright", "Group-scoped export only"),
    ("Writes", "Must carry a contract header; the tag is stamped server-side",
     "Denied"),
    ("Credential allow-list", "Ingest principals", "Payer entitlements"),
]
hx = [0.62, 4.06, 8.4]
hw = [3.3, 4.2, 4.32]
rect(s, 4.06, 1.62, 4.2, 0.5, fill=RGBColor(0x6B, 0x4C, 0x9A),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
rect(s, 8.4, 1.62, 4.32, 0.5, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.14)
txt(s, 4.06, 1.75, 4.2, 0.3, "INBOUND ROUTE", size=11, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER)
txt(s, 8.4, 1.75, 4.32, 0.3, "OUTBOUND ROUTE", size=11, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER)
for j, r in enumerate(rows):
    y = 2.28 + j * 0.72
    if j % 2 == 0:
        rect(s, 0.62, y - 0.08, 12.1, 0.68, fill=WHITE)
    txt(s, 0.72, y + 0.04, 3.1, 0.5, r[0], size=11.5, bold=True, color=GREY)
    txt(s, 4.26, y + 0.04, 3.9, 0.56, r[1], size=11.5, color=INK, line_spacing=1.18)
    txt(s, 8.6, y + 0.04, 4.0, 0.56, r[2], size=11.5, color=INK, line_spacing=1.18)

card(s, 0.62, 5.96, 12.1, 0.86, fill=WHITE)
rect(s, 0.62, 5.96, 0.055, 0.86, fill=AMBER)
txt(s, 1.0, 6.12, 11.5, 0.6,
    "Precisely stated: this is not isolation inside the data store — the rows share a database. What is isolated is the "
    "reachable surface. The control that makes it safe is on the next slide.",
    size=11.5, color=INK, line_spacing=1.24)

# ---------------------------------------------------------------- 10  keystone
s = new(dark=True)
head(s, "The keystone control", "Everything else depends on this one fact", dark=True)
rect(s, 0.62, 1.84, 12.1, 1.9, fill=NAVY_CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.05)
rect(s, 0.62, 1.84, 0.07, 1.9, fill=TEAL_LT)
txt(s, 1.16, 2.22, 11.2, 1.2,
    "No payer application holds any Azure role\non any FHIR service.",
    size=32, bold=True, color=WHITE, font=HEAD, line_spacing=1.14)

for i, (t, d, col) in enumerate([
    ("Authentication succeeds", "Entra issues the payer a perfectly valid token. Nothing is wrong with it.", TEAL_LT),
    ("Authorisation fails", "The FHIR service returns 403, not 401. The identity is known; it simply has no role.", AMBER),
    ("So the gateway cannot be bypassed", "Not by policy or convention — by the absence of an authorisation path.", ICE),
]):
    x = 0.62 + i * 4.1
    rect(s, x, 4.06, 3.86, 1.66, fill=RGBColor(0x14, 0x33, 0x54),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(s, x, 4.06, 3.86, 0.05, fill=col)
    txt(s, x + 0.3, 4.32, 3.3, 0.34, t, size=13.5, bold=True, color=WHITE)
    txt(s, x + 0.3, 4.76, 3.3, 0.86, d, size=11,
        color=RGBColor(0x9F, 0xB8, 0xCC), line_spacing=1.26)

txt(s, 0.62, 6.06, 12.1, 0.6,
    "This is the difference between a boundary and a speed bump — and it is assertion 12 in the proof suite.",
    size=14, bold=True, color=TEAL_LT, font=HEAD)

# ---------------------------------------------------------------- 11  proof
s = new()
head(s, "Proof", "Sixteen assertions. Three minutes. One command.")
rect(s, 0.62, 1.58, 8.5, 0.52, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.1)
txt(s, 0.94, 1.73, 8.0, 0.28, "./scripts/run-isolation-tests.ps1", size=13,
    font=MONO, color=RGBColor(0x8F, 0xE3, 0xC0))
rect(s, 9.32, 1.58, 3.4, 0.52, fill=GREEN_LT, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.1, line=GREEN, lw=0.75)
txt(s, 9.32, 1.72, 3.4, 0.3, "All 16 assertions passed", size=14, bold=True,
    color=GREEN, font=HEAD, align=PP_ALIGN.CENTER)

groups = [
    (0, "The payer can do its job", ["own data readable", "group export accepted",
                                     "capability statement served"], GREEN),
    (0, "Contract scope enforced", ["unentitled export refused",
                                    "caller-supplied filter overridden",
                                    "body carries only own contracts"], TEAL),
    (0, "Direction enforced", ["write refused on read route",
                               "payer refused on ingest route",
                               "export refused on ingest route"], TEAL),
    (1, "Cross-payer access refused", ["valid token, wrong payer",
                                       "wrong audience rejected"], TEAL),
    (1, "Export bounded", ["system-wide export refused", "patient-wide export refused",
                           "second export inside the window throttled"], TEAL),
    (1, "Gateway cannot be bypassed", ["payer token direct to the service refused"], AMBER),
    (1, "Untagged data cannot enter", ["untagged write rejected"], TEAL),
]
cursor = {0: 2.3, 1: 2.3}
for col_i, g, items, col in groups:
    x = 0.62 + col_i * 6.16
    yy = cursor[col_i]
    h = 0.44 + len(items) * 0.28
    card(s, x, yy, 5.94, h)
    rect(s, x, yy, 0.05, h, fill=col)
    txt(s, x + 0.3, yy + 0.12, 5.3, 0.28, g, size=12, bold=True, color=INK)
    for j, it in enumerate(items):
        ry = yy + 0.44 + j * 0.28
        txt(s, x + 0.3, ry, 0.46, 0.24, "PASS", size=8, bold=True, color=GREEN)
        txt(s, x + 0.9, ry - 0.015, 4.86, 0.24, it, size=10.5, color=GREY)
    cursor[col_i] = yy + h + 0.16

txt(s, 0.62, 6.6, 12.1, 0.28,
    "Credentials are minted for the run and revoked at the end. Nothing is left behind.",
    size=10.5, italic=True, color=GREY)

# ---------------------------------------------------------------- 12  reading it
s = new()
head(s, "Reading the proof", "Three results carry the entire argument")
cards = [
    ("04", "A valid token from the wrong payer",
     "Expected 403 · got 403",
     "The token is correctly formed and correctly audienced. It is refused on the entitlement check. "
     "This is the PHI boundary holding even when the credential itself is legitimate — the case that a "
     "naive audience check would miss.", TEAL),
    ("12", "A payer token sent straight to the service",
     "Expected 403 · got 403",
     "Not 401. Entra authenticated the caller; the FHIR service found no role. The gateway is not "
     "optional, and this is the single result that makes every other control meaningful.", AMBER),
    ("13b", "The response body was inspected",
     "Only own contracts returned",
     "Status codes can be satisfied while the payload is wrong. This assertion opens the returned "
     "bundle and checks the contract tags on every resource, not just the HTTP status.", GREEN),
]
for i, (n, t, res, d, col) in enumerate(cards):
    x = 0.62 + i * 4.1
    card(s, x, 1.66, 3.86, 4.28)
    rect(s, x, 1.66, 3.86, 0.055, fill=col)
    txt(s, x + 0.32, 1.96, 1.4, 0.42, n, size=30, bold=True, color=col, font=HEAD)
    txt(s, x + 0.32, 2.5, 3.24, 0.6, t, size=14, bold=True, color=INK,
        line_spacing=1.16)
    rect(s, x + 0.32, 3.2, 3.24, 0.01, fill=LINE)
    txt(s, x + 0.32, 3.36, 3.24, 0.28, res, size=10.5, bold=True, color=GREEN,
        font=MONO)
    txt(s, x + 0.32, 3.8, 3.24, 2.0, d, size=11, color=GREY, line_spacing=1.3)

rect(s, 0.62, 6.16, 12.1, 0.62, fill=ICE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.08)
txt(s, 1.0, 6.34, 11.4, 0.3,
    "Keep the suite in continuous integration. A well-meant role assignment turns assertion 12 red — nothing else notices.",
    size=12.5, bold=True, color=NAVY, font=HEAD)

# ---------------------------------------------------------------- 13  onboarding
s = new()
head(s, "Operations", "Onboarding a payer: seven steps, one command")
rect(s, 0.62, 1.6, 12.1, 0.48, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.1)
txt(s, 0.92, 1.72, 11.6, 0.28,
    "./scripts/onboard-payer.ps1 -PayerKey payerc -DisplayName 'Northwind Health' -Contracts CT-5150",
    size=11.5, font=MONO, color=RGBColor(0x8F, 0xE3, 0xC0))

steps = [
    ("1", "Register the application", "An Entra identity for the payer."),
    ("2", "Issue a credential", "Certificate preferred; the private key never leaves the vault."),
    ("3", "Define the cohort", "A FHIR Group holding the members the contract covers."),
    ("4", "Confirm the data is tagged", "Every resource carries its contract tag, or it is rejected."),
    ("5", "Record the entitlement", "Payer, contracts and permitted groups, in one record."),
    ("6", "Grant no Azure role", "Deliberately. This omission is the control."),
    ("7", "Hand off", "A sheet with endpoints and scope — and no credential material."),
]
for i, (n, t, d) in enumerate(steps):
    x, y = 0.62 + (i % 4) * 3.06, 2.4 + (i // 4) * 1.86
    card(s, x, y, 2.86, 1.66, fill=WHITE,
         line=AMBER if n == "6" else LINE)
    chip(s, x + 0.28, y + 0.26, 0.42, n, fill=AMBER if n == "6" else TEAL, size=13)
    txt(s, x + 0.28, y + 0.84, 2.34, 0.3, t, size=12.5, bold=True, color=INK)
    txt(s, x + 0.28, y + 1.16, 2.34, 0.42, d, size=10, color=GREY,
        line_spacing=1.18)

card(s, 9.8, 4.26, 2.92, 1.66, fill=ICE, line=ICE)
txt(s, 10.08, 4.52, 2.4, 0.3, "OFFBOARDING", size=9.5, bold=True, color=NAVY)
txt(s, 10.08, 4.9, 2.4, 0.86,
    "Delete the Entra application.\nEvery downstream check then fails closed.",
    size=11, color=NAVY, line_spacing=1.24)

rect(s, 0.62, 6.14, 12.1, 0.6, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.07)
txt(s, 1.0, 6.31, 11.4, 0.3,
    "Roughly four minutes, repeatable, and auditable. It is a script, so it cannot drift from the runbook.",
    size=12.5, bold=True, color=WHITE, font=HEAD)

# ---------------------------------------------------------------- 14  the 403
s = new()
head(s, "Unblocked", "The import failure was an identity mismatch, not a permission gap")
ids = [
    ("The operator", "The human or pipeline that starts the job. Needs no storage role at all.", GREY, "NO ROLE NEEDED", GREY),
    ("A user-assigned identity", "The obvious place to put the grant. It is not the identity that reads the blob.", RED, "NOT THIS", RED),
    ("The service's own identity", "The system-assigned identity of the FHIR service. This is what $import actually uses.", GREEN, "USE THIS", GREEN),
]
txt(s, 0.62, 1.62, 6.0, 0.3, "THREE IDENTITIES ARE IN PLAY", size=10, bold=True,
    color=TEAL)
for i, (t, d, col, lab, labcol) in enumerate(ids):
    y = 2.0 + i * 1.24
    card(s, 0.62, y, 6.0, 1.08, fill=WHITE, line=GREEN if labcol == GREEN else LINE)
    rect(s, 0.62, y, 0.055, 1.08, fill=col)
    txt(s, 1.0, y + 0.2, 4.2, 0.3, t, size=13, bold=True, color=INK)
    txt(s, 1.0, y + 0.52, 4.2, 0.5, d, size=10.5, color=GREY, line_spacing=1.2)
    txt(s, 5.2, y + 0.34, 1.3, 0.3, lab, size=9, bold=True, color=labcol,
        align=PP_ALIGN.RIGHT)

card(s, 6.9, 1.94, 5.82, 1.7, fill=WHITE)
rect(s, 6.9, 1.94, 0.055, 1.7, fill=TEAL)
txt(s, 7.28, 2.16, 5.1, 0.3, "THE FIX", size=10, bold=True, color=TEAL)
txt(s, 7.28, 2.5, 5.1, 1.0,
    "Grant Storage Blob Data Contributor to each FHIR service's own system-assigned "
    "principal. It is one module in the template, so it cannot be forgotten on the "
    "fortieth payer.", size=11.5, color=INK, line_spacing=1.28)

card(s, 6.9, 3.82, 5.82, 2.02, fill=AMBER_LT, line=AMBER)
txt(s, 7.28, 4.04, 5.1, 0.3, "AND WHY THE RETEST LOOKED WRONG", size=10, bold=True,
    color=AMBER)
txt(s, 7.28, 4.38, 5.1, 1.3,
    "A completed import job keeps its result. Re-polling the old job replays the "
    "original 403 byte-for-byte, even after the grant is correct.\n\n"
    "Storage logs settle it: no request logged means the 403 is a replay, not a refusal.",
    size=11, color=INK, line_spacing=1.28)

txt(s, 0.62, 6.1, 12.1, 0.6,
    "A 403 that looks identical whether the grant is missing or the answer is cached is how a one-line fix costs a week. "
    "The runbook now separates the two in a single query.",
    size=11.5, color=GREY, line_spacing=1.3)

# ---------------------------------------------------------------- 15  observability
s = new()
head(s, "Evidence", "Enforced, and visible")
obs = [
    ("Gateway denials", "Every 401, 403 and 429, grouped by route and payer.",
     "Shows the isolation model working rather than merely configured — the proof "
     "run appears here within a minute.", TEAL),
    ("Storage refusals", "Every 403 with the object id of the identity that was refused.",
     "Turns 'permissions are wrong somewhere' into a named principal. This is the "
     "query that separates a live refusal from a replayed one.", AMBER),
    ("Export behaviour", "Job durations, queue depth and throttle responses per payer.",
     "The baseline the capacity work will be measured against, collecting from "
     "day one rather than from the day it is needed.", GREEN),
]
for i, (t, what, why, col) in enumerate(obs):
    x = 0.62 + i * 4.1
    card(s, x, 1.66, 3.86, 3.4)
    rect(s, x, 1.66, 3.86, 0.055, fill=col)
    txt(s, x + 0.32, 1.96, 3.24, 0.34, t, size=15, bold=True, color=INK, font=HEAD)
    txt(s, x + 0.32, 2.4, 3.24, 0.66, what, size=11.5, color=INK, line_spacing=1.24)
    rect(s, x + 0.32, 3.16, 3.24, 0.01, fill=LINE)
    txt(s, x + 0.32, 3.34, 3.24, 1.5, why, size=10.5, color=GREY, line_spacing=1.3)

rect(s, 0.62, 5.24, 12.1, 1.44, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.06)
txt(s, 1.0, 5.48, 11.4, 0.28, "PINNED FROM DAY ONE", size=10, bold=True,
    color=TEAL_LT)
txt(s, 1.0, 5.82, 11.4, 0.66,
    "Three saved queries ship with the environment. An auditor asking \"show me that payer B has never "
    "reached payer A's data\" gets an answer from data, in about a minute — not from a diagram.",
    size=13, color=ICE, line_spacing=1.3)

# ---------------------------------------------------------------- 16  the gap
s = new()
head(s, "Open item", "One question we will not answer from documentation")
card(s, 0.62, 1.66, 12.1, 1.18, fill=AMBER_LT, line=AMBER)
rect(s, 0.62, 1.66, 0.06, 1.18, fill=AMBER)
txt(s, 1.06, 1.9, 11.3, 0.34,
    "Concurrent bulk-export capacity is unpublished, and therefore unknown.",
    size=17, bold=True, color=INK, font=HEAD)
txt(s, 1.06, 2.3, 11.3, 0.4,
    "Roughly forty payers pulling on the same schedule is the stated concern. No documentation states how many "
    "simultaneous export jobs a service sustains, or how quickly it scales into a burst.",
    size=11.5, color=GREY, line_spacing=1.24)

txt(s, 0.62, 3.14, 5.9, 0.3, "DEPLOYED TODAY — REDUCES THE RISK", size=10,
    bold=True, color=TEAL)
for i, t in enumerate(["One concurrent export per payer per five-minute window",
                       "Cohort-scoped export only — never system-wide or per-patient",
                       "600 requests per minute and 50,000 per day, per payer"]):
    y = 3.5 + i * 0.72
    card(s, 0.62, y, 5.9, 0.6)
    rect(s, 0.62, y, 0.05, 0.6, fill=TEAL)
    txt(s, 1.0, y + 0.16, 5.3, 0.36, t, size=11.5, color=INK)

txt(s, 6.82, 3.14, 5.9, 0.3, "NOT YET DONE — MEASURES IT", size=10, bold=True,
    color=AMBER)
card(s, 6.82, 3.5, 5.9, 1.82, fill=WHITE, line=AMBER)
txt(s, 7.2, 3.72, 5.2, 1.5,
    "The load harness is written. It drives concurrency 1 → 5 → 15 → 40 with an "
    "import running in parallel, and reports p95, median and throttle rate for each "
    "phase.\n\nWhat it produces is a number, and a defensible ceiling per payer.",
    size=11.5, color=INK, line_spacing=1.3)

rect(s, 6.82, 5.5, 5.9, 0.72, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.07)
txt(s, 7.2, 5.7, 5.2, 0.34,
    "Run it before a go-live date is committed, not after.",
    size=13, bold=True, color=WHITE, font=HEAD)

txt(s, 0.62, 6.42, 5.9, 0.34,
    "Mitigations reduce risk. They do not measure it.",
    size=12, bold=True, italic=True, color=GREY)

# ---------------------------------------------------------------- 17  cost
s = new()
head(s, "Cost", "Forty instances is not forty times the cost")
for i, (v, l, col) in enumerate([("~$450", "per month, POC running continuously", TEAL),
                                 ("~$10", "per demo day, deleted in between", GREEN),
                                 ("~$2,200", "per month, production estimate at scale", NAVY)]):
    x = 0.62 + i * 4.1
    card(s, x, 1.66, 3.86, 1.72)
    rect(s, x, 1.66, 3.86, 0.055, fill=col)
    stat(s, x + 0.32, 1.94, 3.2, v, l, color=col, vsize=32)

card(s, 0.62, 3.6, 5.94, 2.42)
rect(s, 0.62, 3.6, 0.055, 2.42, fill=TEAL)
txt(s, 1.0, 3.86, 5.2, 0.3, "WHY THE MULTIPLE DOES NOT APPLY", size=10, bold=True,
    color=TEAL)
for i, t in enumerate([
        "FHIR services bill on provisioned throughput and stored data — not per instance.",
        "Splitting the same workload across more services divides the throughput; it does not duplicate it.",
        "The gateway, Log Analytics and storage are shared fixed costs across all payers.",
        "The variable cost that actually scales is stored data, and that is the same volume either way."]):
    y = 4.24 + i * 0.44
    chip(s, 1.0, y + 0.03, 0.16, "", fill=TEAL)
    txt(s, 1.3, y, 4.9, 0.42, t, size=10.5, color=INK, line_spacing=1.18)

card(s, 6.78, 3.6, 5.94, 2.42, fill=WHITE)
rect(s, 6.78, 3.6, 0.055, 2.42, fill=AMBER)
txt(s, 7.16, 3.86, 5.2, 0.3, "THE ONE COST CONTROL THAT MATTERS", size=10,
    bold=True, color=AMBER)
txt(s, 7.16, 4.24, 5.2, 1.5,
    "A FHIR service has no pause state. A provisioned service bills whether or not "
    "anyone calls it.\n\nFor a POC, the environment rebuilds from the template in "
    "about twenty minutes. Deleting it between sessions is the difference between "
    "$450 a month and $10 a day.",
    size=11.5, color=INK, line_spacing=1.3)

txt(s, 0.62, 6.24, 12.1, 0.34,
    "Cost was never the reason to reject per-payer isolation. Instance-count management was — and a template answers that.",
    size=12, bold=True, color=NAVY, font=HEAD)

# ---------------------------------------------------------------- 18  traceability
s = new()
head(s, "Traceability", "Every ask, and the artefact that answers it")
tr = [
    ("Reference architecture", "Delivered", "Editable four-page architecture set", GREEN),
    ("Physical vs logical isolation", "Deployed", "Two services; assertions 4 and 13b", GREEN),
    ("PHI separation at the payer", "Deployed", "Assertion 4 — valid token, wrong payer, refused", GREEN),
    ("Payers kept out of ingest", "Deployed", "Assertions 7, 8 and 9", GREEN),
    ("Gateway is mandatory", "Deployed", "Six-layer policy; assertion 12", GREEN),
    ("Payer onboarding process", "Scripted", "Seven steps, roughly four minutes", GREEN),
    ("The $import 403", "Root-caused and fixed", "Grant corrected in the template", GREEN),
    ("Quota 10 to 40", "Answered", "Support request; treat as a lead-time item", GREEN),
    ("Concurrent export capacity", "Open — mitigated", "Harness written; measurement pending", AMBER),
]
rect(s, 0.62, 1.62, 12.1, 0.42, fill=NAVY)
for hx_, hw_, lab in [(0.92, 4.3, "THE ASK"), (5.3, 2.5, "STATUS"),
                      (8.0, 4.5, "WHERE IT IS ANSWERED")]:
    txt(s, hx_, 1.72, hw_, 0.26, lab, size=9.5, bold=True, color=TEAL_LT)
for i, (a, st, ev, col) in enumerate(tr):
    y = 2.12 + i * 0.5
    if i % 2 == 0:
        rect(s, 0.62, y - 0.04, 12.1, 0.48, fill=WHITE)
    txt(s, 0.92, y + 0.06, 4.2, 0.3, a, size=11.5, bold=True, color=INK)
    txt(s, 5.3, y + 0.06, 2.5, 0.3, st, size=11, bold=True, color=col)
    txt(s, 8.0, y + 0.06, 4.5, 0.3, ev, size=11, color=GREY)

rect(s, 0.62, 6.68, 12.1, 0.02, fill=LINE)

# ---------------------------------------------------------------- 19  proven / next
s = new()
head(s, "Scope", "What is proven, and what production still needs")
card(s, 0.62, 1.66, 5.94, 4.86)
rect(s, 0.62, 1.66, 5.94, 0.055, fill=GREEN)
txt(s, 1.0, 1.96, 5.2, 0.34, "PROVEN IN THE DEPLOYED ENVIRONMENT", size=10,
    bold=True, color=GREEN)
for i, t in enumerate(["Payer isolation, tested rather than asserted",
                       "Contract scoping enforced on every read",
                       "Inbound and outbound directions fully separated",
                       "The gateway cannot be bypassed",
                       "Onboarding and offboarding, scripted",
                       "The import failure, root-caused and fixed",
                       "Denials and refusals, observable in queries"]):
    y = 2.42 + i * 0.56
    chip(s, 1.0, y + 0.04, 0.2, "", fill=GREEN)
    txt(s, 1.36, y, 4.9, 0.42, t, size=12, color=INK, line_spacing=1.18)

card(s, 6.78, 1.66, 5.94, 4.86)
rect(s, 6.78, 1.66, 5.94, 0.055, fill=AMBER)
txt(s, 7.16, 1.96, 5.2, 0.34, "REQUIRED BEFORE PRODUCTION", size=10, bold=True,
    color=AMBER)
nxt = [("Private endpoints", "Low"), ("Certificate credentials, not secrets", "Low"),
       ("Entitlements sourced from the contract master", "Medium"),
       ("Customer-managed keys", "Low"),
       ("Quarantine handling for rejected inbound", "Medium"),
       ("Multi-region continuity", "Medium"),
       ("A measured capacity envelope", "Low")]
for i, (t, e) in enumerate(nxt):
    y = 2.42 + i * 0.56
    chip(s, 7.16, y + 0.04, 0.2, "", fill=AMBER)
    txt(s, 7.52, y, 4.0, 0.42, t, size=12, color=INK, line_spacing=1.18)
    txt(s, 11.6, y, 0.9, 0.3, e, size=9.5, bold=True, color=GREY,
        align=PP_ALIGN.RIGHT)

txt(s, 0.62, 6.62, 12.1, 0.3,
    "None of the production items changes the architecture. They are configuration and hardening on top of it.",
    size=11.5, italic=True, color=GREY)

# ---------------------------------------------------------------- 20  next steps
s = new(dark=True)
head(s, "Recommendation", "Four things, in this order", dark=True)
acts = [
    ("01", "Run the capacity harness", "Before any production date is committed. It is the only open question, and the instrument already exists.", "This week"),
    ("02", "File the quota request", "Ten to forty FHIR services is a support request evaluated against regional capacity. It can be declined. It is a lead-time item.", "This week"),
    ("03", "Move entitlements to the contract master", "Today they live in gateway configuration, which is fine for a proof and wrong for production.", "Next iteration"),
    ("04", "Put the proof suite in the pipeline", "It is the only thing that notices when a role assignment quietly turns a boundary into a suggestion.", "Next iteration"),
]
for i, (n, t, d, when) in enumerate(acts):
    y = 1.78 + i * 1.13
    rect(s, 0.62, y, 12.1, 0.98, fill=NAVY_CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.06)
    rect(s, 0.62, y, 0.05, 0.98, fill=TEAL_LT)
    txt(s, 1.02, y + 0.28, 0.6, 0.4, n, size=17, bold=True, color=TEAL_LT, font=HEAD)
    txt(s, 1.86, y + 0.16, 4.3, 0.32, t, size=14, bold=True, color=WHITE)
    txt(s, 1.86, y + 0.52, 8.4, 0.4, d, size=10.5,
        color=RGBColor(0x9F, 0xB8, 0xCC), line_spacing=1.2)
    txt(s, 10.9, y + 0.34, 1.6, 0.3, when, size=10.5, bold=True, color=ICE,
        align=PP_ALIGN.RIGHT)

txt(s, 0.62, 6.34, 12.1, 0.42,
    "The environment is live. The proof runs in three minutes. The remaining work is measurement and hardening, not design.",
    size=14, bold=True, color=TEAL_LT, font=HEAD)

out = Path(__file__).with_name("AHDS-Payer-Isolation-Reference-Implementation.pptx")
prs.core_properties.title = "Payer Data Isolation on Azure Health Data Services"
prs.core_properties.author = "Microsoft Health & Life Sciences"
prs.save(out)
print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")
