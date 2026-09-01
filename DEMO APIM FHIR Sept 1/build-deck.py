"""Northwind Health - APIM + AHDS session deck, 2026-09-01.

Structured to the eight agenda items, one section each, speaker notes written to
be read aloud verbatim.

    %LOCALAPPDATA%\\venvs\\pptxbuild\\Scripts\\python.exe build-deck.py

Facts are from the live environment in rg-ahds-fhir-poc, verified 2026-08-31.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

W, H = 13.333, 7.5

NAVY = RGBColor(0x0E, 0x2A, 0x47)
NAVY_CARD = RGBColor(0x1A, 0x3E, 0x60)
TEAL = RGBColor(0x1C, 0x72, 0x93)
TEAL_LT = RGBColor(0x4F, 0xA6, 0xC4)
ICE = RGBColor(0xCF, 0xE3, 0xEF)
PAPER = RGBColor(0xF5, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x12, 0x28, 0x3D)
GREY = RGBColor(0x5A, 0x6B, 0x7B)
LINE = RGBColor(0xD8, 0xE2, 0xEA)
AMBER = RGBColor(0xD9, 0x82, 0x2B)
AMBER_LT = RGBColor(0xFD, 0xF3, 0xE3)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
GREEN_LT = RGBColor(0xE8, 0xF4, 0xEC)
RED = RGBColor(0xB0, 0x3A, 0x2B)
RED_LT = RGBColor(0xFA, 0xEC, 0xEA)
SLATE = RGBColor(0x2C, 0x50, 0x74)
MUTED = RGBColor(0x9F, 0xB8, 0xCC)

HEAD = "Cambria"
BODY = "Calibri"
MONO = "Consolas"

GATEWAY = "https://apim-poc-ahds-demo01.azure-api.net"

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

_n = {"i": 0}


def new(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY if dark else PAPER
    _n["i"] += 1
    return s


def rect(s, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE, radius=None, line=None, lw=1.0):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def txt(s, x, y, w, h, runs, size=14, color=INK, font=BODY, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, italic=False,
        line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    if isinstance(runs, str):
        runs = [(runs, {})]
    runs = [(r, {}) if isinstance(r, str) else r for r in runs]

    for i, (text, ov) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("line_spacing", line_spacing)
        if i > 0:
            p.space_before = Pt(ov.get("space_before", space))
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic)
        f.color.rgb = ov.get("color", color)
    return tb


def foot(s, dark=False):
    rect(s, 0.62, 6.94, 12.1, 0.014, fill=SLATE if dark else LINE)
    txt(s, 0.62, 7.03, 9.6, 0.24,
        "Northwind Health  \u00b7  APIM + Azure Health Data Services  \u00b7  1 September 2026",
        size=9, color=RGBColor(0x7F, 0x9A, 0xB5) if dark else GREY)
    txt(s, 11.0, 7.03, 1.72, 0.24, f"{_n['i']:02d}", size=9, align=PP_ALIGN.RIGHT,
        color=RGBColor(0x7F, 0x9A, 0xB5) if dark else GREY)


def head(s, kicker, title, dark=False):
    rect(s, 0.62, 0.46, 0.075, 0.72, fill=TEAL_LT if dark else TEAL)
    txt(s, 0.85, 0.46, 11.8, 0.26, kicker.upper(), size=10.5, bold=True,
        color=TEAL_LT if dark else TEAL, font=BODY)
    txt(s, 0.85, 0.72, 11.8, 0.52, title, size=27, bold=True,
        color=WHITE if dark else INK, font=HEAD)
    foot(s, dark)


def card(s, x, y, w, h, fill=WHITE, line=LINE, radius=0.06):
    return rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                radius=radius, line=line, lw=0.75)


def chip(s, x, y, d, label, fill=TEAL, color=WHITE, size=12):
    rect(s, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)
    txt(s, x, y + d / 2 - 0.115, d, 0.25, label, size=size, bold=True,
        color=color, align=PP_ALIGN.CENTER, font=BODY)


def code(s, x, y, w, h, lines, size=11.5, fill=RGBColor(0x0B, 0x22, 0x3A)):
    rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    runs = []
    for i, ln in enumerate(lines):
        if isinstance(ln, tuple):
            text, col = ln
        else:
            text, col = ln, ICE
        runs.append((text, {"color": col, "space_before": 0 if i == 0 else 3}))
    txt(s, x + 0.24, y + 0.2, w - 0.48, h - 0.4, runs, size=size, font=MONO,
        line_spacing=1.18)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


def bullets(s, x, y, w, items, size=14, gap=9, color=INK, bold_lead=True):
    """items = list of (lead, rest). Lead is bolded, rest follows on the line."""
    runs = []
    for i, (lead, rest) in enumerate(items):
        runs.append((lead, {"bold": bold_lead, "color": color, "size": size,
                            "space_before": 0 if i == 0 else gap}))
        if rest:
            runs.append((rest, {"color": GREY, "size": size - 0.5,
                                "space_before": 2}))
    txt(s, x, y, w, 4.4, runs, size=size, line_spacing=1.14)


# =====================================================================  01 title
s = new(dark=True)
rect(s, 0, 0, W, 0.09, fill=TEAL)
txt(s, 0.95, 1.72, 11.4, 0.34, "NORTHWIND HEALTH  \u00b7  CMS-0057-F",
    size=12.5, bold=True, color=TEAL_LT)
txt(s, 0.95, 2.16, 11.4, 1.5,
    "Azure API Management in front of\nAzure Health Data Services",
    size=41, bold=True, color=WHITE, font=HEAD, line_spacing=1.06)
rect(s, 0.95, 3.86, 2.0, 0.03, fill=TEAL)
txt(s, 0.95, 4.12, 11.0, 0.9,
    "How payer isolation is enforced, proven with a live gateway, "
    "running policies and a test suite whose assertions are the security guarantees.",
    size=15.5, color=ICE, line_spacing=1.3)
txt(s, 0.95, 5.72, 11.4, 0.8,
    [("Working session  \u00b7  1 September 2026", {"size": 13, "color": MUTED}),
     ("Microsoft Health & Life Sciences  \u00b7  Enterprise Architecture, Integration Engineering, Platform Engineering",
      {"size": 12, "color": SLATE, "space_before": 4})],
    size=13)
foot(s, dark=True)
notes(s, """
Open here. Do not read the title.

Say: "Thanks for the time. Last time we met, on the twelfth, we drew four options on
the whiteboard and picked one. What I have today is that option, deployed. Not slides
about it - a running gateway, the actual policies, and a test suite you can read.

The session has eight parts. I want to spend most of it on two of them: how inbound and
outbound are actually kept apart, and the tests that prove a payer cannot reach data that
isn't theirs. Everything else is context for those two.

One thing before we start - please interrupt. If something doesn't hold up, I would much
rather find out in this room than in your architecture review."

Then move straight to the agenda slide.
""")

# ====================================================================  02 agenda
s = new()
head(s, "Agenda", "Eight parts, sixty minutes")
items = [
    ("1", "Objectives and scope", "What we are answering today, and what we are not", "05"),
    ("2", "Overview of APIM", "What the product is, for those who have not used it", "07"),
    ("3", "APIM in the CMS architecture", "Where it sits, and what it is responsible for", "07"),
    ("4", "Inbound vs outbound partitioning", "The model, and why one FHIR service can serve both", "08"),
    ("5", "Gateway isolation with tests", "Live. Sixteen assertions against the deployed gateway", "12"),
    ("6", "Policy walkthrough and debug trace", "The six layers, then a real trace of one call", "10"),
    ("7", "Credentials, identity and deployment", "How a payer authenticates, and what you run in prod", "07"),
    ("8", "Decisions and next steps", "Four decisions I need from this room", "04"),
]
y = 1.62
for num, title, sub, mins in items:
    hl = num in ("4", "5", "6")
    card(s, 0.62, y, 12.1, 0.6, fill=WHITE if not hl else RGBColor(0xEC, 0xF4, 0xF8),
         line=LINE if not hl else TEAL_LT)
    chip(s, 0.82, y + 0.14, 0.32, num, fill=TEAL if hl else RGBColor(0xB9, 0xC9, 0xD6), size=11)
    txt(s, 1.32, y + 0.11, 4.6, 0.28, title, size=14, bold=True, color=INK)
    txt(s, 1.32, y + 0.34, 8.6, 0.24, sub, size=11, color=GREY)
    txt(s, 11.55, y + 0.19, 0.95, 0.26, f"{mins} min", size=11, bold=True,
        color=TEAL if hl else GREY, align=PP_ALIGN.RIGHT)
    y += 0.665
notes(s, """
Say: "Here is the shape of the hour. The three highlighted rows - four, five and six - are
where the substance is. That is about thirty minutes in the middle where I stop talking
and start showing you the running system.

Parts one to three are set-up. If you already know API Management, part two will be
familiar and I will move quickly. Part seven is the operational reality - credentials and
what you would actually deploy. Part eight is decisions, and I do need those today because
they gate the next block of work."

If someone asks to reorder: the demo needs parts three and four as set-up, so keep the
order, but you can compress two.

Timing discipline: if you are behind at the half hour, cut part six down to the debug
trace only. Never cut part five.
""")

# =========================================================  03 objectives + scope
s = new()
head(s, "1 \u00b7 Objectives and scope", "What this session has to settle")
card(s, 0.62, 1.6, 5.95, 3.62, fill=WHITE)
txt(s, 0.92, 1.85, 5.4, 0.3, "IN SCOPE - we answer these today", size=11, bold=True, color=TEAL)
bullets(s, 0.92, 2.25, 5.4, [
    ("Can one FHIR service serve inbound and outbound safely?",
     "Platform Engineering's question from 8/12, answered with running policy"),
    ("What stops payer A reading payer B?",
     "Demonstrated, not asserted - assertion 4 and assertion 12"),
    ("How do contracts partition inside one payer?",
     "meta.tag plus enforced _tag injection at the gateway"),
    ("What does onboarding payer number three cost us?",
     "A configuration change - shown live in part four"),
], size=12.5, gap=10)

card(s, 6.78, 1.6, 5.94, 3.62, fill=RGBColor(0xF7, 0xF9, 0xFB))
txt(s, 7.08, 1.85, 5.4, 0.3, "OUT OF SCOPE - flagged, not solved", size=11, bold=True, color=AMBER)
bullets(s, 7.08, 2.25, 5.4, [
    ("Production capacity sizing",
     "Needs your real export volumes; the model is in docs/05"),
    ("SMART on FHIR scope-level filtering",
     "AHDS does not enforce contract scopes - that is why APIM is here"),
    ("Payer contracting and legal",
     "Out of engineering's hands"),
    ("Migration of the existing dev workspace",
     "Separate workstream once the pattern is agreed"),
], size=12.5, gap=10)

card(s, 0.62, 5.55, 12.1, 1.12, fill=GREEN_LT, line=GREEN)
txt(s, 0.92, 5.76, 11.5, 0.74,
    [("The single question behind all of it:  ", {"bold": True, "size": 14, "color": INK}),
     ("when two payers call the same platform, what physically prevents one from reading the other's members?",
      {"size": 14, "color": INK})], size=14, line_spacing=1.22)
notes(s, """
Say: "Objectives first, so we all agree on what a good outcome looks like.

On the left, four things I intend to have settled by the time we leave. These are your
questions, not mine - three of them came off the eight-twelve transcript and one came from
Integration Engineering by email.

On the right, four things I am deliberately not solving today. I want to be honest about
these rather than have them surface as gotchas. Capacity sizing is the big one - I cannot
size your production gateway until you give me real export volumes, and I will come back
to that in part seven.

And the line at the bottom is the question everything else serves. If you take one thing
away, take that sentence. When two payers call the same platform, what physically prevents
one from reading the other's members? Not 'what policy says they shouldn't' - what
physically prevents it."

Pause after reading the green bar. Let it sit. That is the thesis of the session.
""")

# ==========================================================  04 where we left off
s = new()
head(s, "1 \u00b7 Objectives and scope", "Where we left off on 12 August")
txt(s, 0.62, 1.55, 12.1, 0.32,
    "Two quotes from the working session. Everything built since is a response to them.",
    size=13, color=GREY)

card(s, 0.62, 2.0, 5.95, 1.9, fill=WHITE, line=TEAL_LT)
txt(s, 0.92, 2.22, 5.35, 1.0,
    "\u201cSeparate instances for payer \u2026 and then within a payer, if we want to separate "
    "by different contracts, that\u2019s going to be a logical separation instead of a physical one.\u201d",
    size=13, italic=True, color=INK, line_spacing=1.22)
txt(s, 0.92, 3.42, 5.35, 0.28, "Enterprise Architecture  \u00b7  40:48", size=11, bold=True, color=TEAL)

card(s, 6.78, 2.0, 5.94, 1.9, fill=WHITE, line=TEAL_LT)
txt(s, 7.08, 2.22, 5.35, 1.0,
    "\u201cFor inbound and outbound data, can it be isolated from each other within one service? "
    "\u2026 we don\u2019t want payer to query our inbound data.\u201d",
    size=13, italic=True, color=INK, line_spacing=1.22)
txt(s, 7.08, 3.42, 5.35, 0.28, "Platform Engineering  \u00b7  8 August working session", size=11, bold=True, color=TEAL)

txt(s, 0.62, 4.12, 12.1, 0.3, "WHAT CHANGED SINCE", size=11, bold=True, color=TEAL)
rows = [
    ("Reference architecture", "Delivered - four pages, editable draw.io"),
    ("The isolation model", "Deployed - two payers, three contracts, running now"),
    ("The $import 403", "Root-caused and fixed in the Bicep, annotated so it stays fixed"),
    ("Inbound / outbound split", "Built as two APIM APIs with mirrored allow-lists"),
]
y = 4.5
for a, b in rows:
    card(s, 0.62, y, 12.1, 0.5, fill=WHITE)
    txt(s, 0.92, y + 0.14, 3.5, 0.26, a, size=12.5, bold=True, color=INK)
    txt(s, 4.5, y + 0.14, 8.0, 0.26, b, size=12.5, color=GREY)
    y += 0.56
notes(s, """
Say: "I want to anchor this in your words, not mine.

The first quote is the architecture decision this room made. Separate instances per payer,
and inside a payer, contracts are a logical separation. That is not my design - it is the
one you landed on, and everything I built follows it.

The second quote is the question that was still open when we ran out of time. Can inbound
and outbound be isolated inside one service? The concern was specific and it was a good
one - you do not want a payer reaching back into the claims data they just sent you.

The four rows at the bottom are what has changed since. The architecture diagram was asked
for twice on that call and I owed it to you. The import 403 that Platform Engineering hit -
I found the root cause, and I will show you the one-line version in part seven if we have
time. The full write-up is in the docs folder."

If asked about the $import 403 now: "It was reading storage as the FHIR service's own
system identity, not the user-assigned identity that held the role. One line of Bicep.
Happy to go deep after, it is written up in docs slash zero-four."
""")

# ====================================================  05 what APIM is (overview)
s = new()
head(s, "2 \u00b7 Overview of APIM", "What API Management actually is")
txt(s, 0.62, 1.52, 12.1, 0.3,
    "A reverse proxy you can program. Every call to your API goes through it, and you decide what happens on the way.",
    size=13.5, color=GREY)

comps = [
    ("Gateway", "The runtime. Receives every call, applies policy, forwards to the backend.", TEAL),
    ("API", "A route definition. We have four - two per payer, inbound and outbound.", TEAL),
    ("Policy", "XML executed per call. Auth, rewriting, filtering, limits. The logic lives here.", AMBER),
    ("Named value", "Configuration and secrets, editable without redeploying. Our entitlement map.", TEAL),
    ("Managed identity", "The gateway's own Entra identity. This is what calls FHIR - not the payer.", GREEN),
    ("Product / subscription", "Optional packaging and keys. We do not use them - Entra tokens instead.", GREY),
]
y = 2.0
for i, (name, desc, col) in enumerate(comps):
    x = 0.62 if i % 2 == 0 else 6.78
    if i % 2 == 0 and i > 0:
        y += 0.92
    card(s, x, y, 5.95, 0.8, fill=WHITE)
    rect(s, x, y, 0.055, 0.8, fill=col)
    txt(s, x + 0.28, y + 0.13, 5.4, 0.26, name, size=13.5, bold=True, color=INK)
    txt(s, x + 0.28, y + 0.42, 5.5, 0.3, desc, size=11.5, color=GREY, line_spacing=1.1)

card(s, 0.62, 4.92, 12.1, 1.3, fill=AMBER_LT, line=AMBER)
txt(s, 0.92, 5.12, 11.5, 0.95,
    [("Why this matters here:  ", {"bold": True, "size": 13.5, "color": INK}),
     ("AHDS role-based access control has no concept of \u201ccontract\u201d. It can say \u201cthis identity may read this FHIR service\u201d "
      "and nothing finer. Every rule finer than that has to be enforced somewhere \u2014 and this is the somewhere.",
      {"size": 13.5, "color": INK})], size=13.5, line_spacing=1.24)
notes(s, """
Say: "Quick level-set, because I do not want to assume everyone here has used API
Management. If you have, give me ninety seconds and we move on.

APIM is a reverse proxy that you can program. Every call to your API goes through it, and
on the way through you get to run logic.

Six pieces. The gateway is the runtime. An API is a route - we have four, two per payer.
A policy is XML that runs on every call, and that is where all the interesting logic lives.
Named values are configuration you can change without redeploying - our payer entitlement
map is one. The managed identity is the gateway's own identity in Entra, and I have flagged
it green because it is the single most important thing on this slide - I will come back to
it in part six. Products and subscriptions we are not using; we authenticate with Entra
tokens instead of API keys, because that is what SMART Backend Services expects."

Then the amber bar - read it close to verbatim:

"And here is why any of this matters. AHDS role-based access control has no concept of a
contract. It can say 'this identity may read this FHIR service' and nothing finer than
that. So every rule finer than the whole service has to be enforced somewhere else. This
is the somewhere else."
""")

# ==========================================================  06 the policy pipeline
s = new()
head(s, "2 \u00b7 Overview of APIM", "The policy pipeline \u2014 four stages per call")
stages = [
    ("inbound", "Before the backend is called", "Validate the token. Resolve entitlement. Reject, rewrite or\nenrich the request. Almost everything we do happens here.", TEAL, "378 lines"),
    ("backend", "The forwarding step", "Attach the managed identity token and call AHDS.\nThe payer's own token is never forwarded.", GREEN, "layer 6"),
    ("outbound", "After the backend responds", "Rewrite response headers so the payer never learns the\nAHDS hostname. Strip internal detail.", TEAL, "Content-Location"),
    ("on-error", "If anything throws", "Return a FHIR OperationOutcome rather than an APIM\nstack trace. Errors stay in the payer's own vocabulary.", RED, "OperationOutcome"),
]
y = 1.72
for name, when, what, col, tag in stages:
    card(s, 0.62, y, 12.1, 1.16, fill=WHITE)
    rect(s, 0.62, y, 0.055, 1.16, fill=col)
    txt(s, 0.95, y + 0.16, 2.5, 0.3, f"<{name}>", size=15, bold=True, color=col, font=MONO)
    txt(s, 0.95, y + 0.52, 2.6, 0.28, when, size=11, color=GREY)
    txt(s, 3.85, y + 0.2, 7.0, 0.8, what, size=12.5, color=INK, line_spacing=1.2)
    txt(s, 11.1, y + 0.44, 1.45, 0.28, tag, size=10.5, bold=True, color=col, align=PP_ALIGN.RIGHT)
    y += 1.27
notes(s, """
Say: "One more concept and then we get to your architecture.

A policy has four stages and they run in this order on every single call.

Inbound runs before your backend is touched. That is where we validate the token, work out
which contracts the caller is entitled to, and reject the call if anything is wrong. Almost
all of our logic is here - three hundred and seventy-eight lines of it.

Backend is the forwarding step. This is where the gateway attaches its own managed identity
token and calls AHDS. I want to be precise about this: the payer's token is validated and
then discarded. It is never forwarded to FHIR. The call that reaches AHDS is made by the
gateway, as the gateway.

Outbound runs on the way back. We rewrite response headers - specifically Content-Location
on a bulk export, so the payer gets a gateway URL to poll and never learns the AHDS
hostname.

And on-error means that when something fails, the payer gets a FHIR OperationOutcome - a
proper FHIR error resource - instead of an API Management stack trace. Small thing, but it
means the payer's own FHIR client can parse our failures."

If asked "is that a lot of policy?": "It is about five hundred and sixty lines across both
routes, heavily commented. The commented version is what you will see in part six."
""")

# =========================================================  07 APIM in the CMS arch
s = new()
head(s, "3 \u00b7 APIM in the CMS architecture", "Where the gateway sits")

txt(s, 0.62, 1.5, 12.1, 0.28, "PAYERS", size=10.5, bold=True, color=GREY)
card(s, 0.62, 1.8, 2.85, 1.0, fill=WHITE, line=TEAL_LT)
txt(s, 0.82, 1.98, 2.5, 0.26, "Contoso Health Plan", size=12.5, bold=True, color=INK)
txt(s, 0.82, 2.26, 2.5, 0.44, "CT-3456, CT-7788\ncmsdqm-payera", size=10.5, color=GREY, line_spacing=1.15)
card(s, 3.62, 1.8, 2.85, 1.0, fill=WHITE, line=TEAL_LT)
txt(s, 3.82, 1.98, 2.5, 0.26, "Fabrikam Medicare Adv.", size=12.5, bold=True, color=INK)
txt(s, 3.82, 2.26, 2.5, 0.44, "CT-9001\ncmsdqm-payerb", size=10.5, color=GREY, line_spacing=1.15)

card(s, 7.0, 1.8, 5.72, 1.0, fill=NAVY, line=NAVY)
txt(s, 7.25, 1.96, 5.2, 0.26, "Microsoft Entra ID", size=12.5, bold=True, color=WHITE)
txt(s, 7.25, 2.26, 5.2, 0.44,
    "client_credentials  \u00b7  no user present\nSMART on FHIR Backend Services shape",
    size=10.5, color=ICE, line_spacing=1.15)

rect(s, 0.62, 3.0, 12.1, 0.02, fill=LINE)

txt(s, 0.62, 3.14, 12.1, 0.28, "GATEWAY  \u00b7  the only path in", size=10.5, bold=True, color=TEAL)
card(s, 0.62, 3.44, 12.1, 1.42, fill=RGBColor(0xEC, 0xF4, 0xF8), line=TEAL)
txt(s, 0.92, 3.6, 4.2, 0.28, "Azure API Management", size=14, bold=True, color=INK)
txt(s, 0.92, 3.9, 4.4, 0.6, "apim-poc-ahds-demo01\nBasicV2  \u00b7  East US 2", size=11, color=GREY, font=MONO, line_spacing=1.2)
layers = ["1 validate-jwt", "2 entitlement", "3 route allow-list",
          "4 contract scoping", "5 rate limit", "6 managed identity"]
x = 5.5
for i, l in enumerate(layers):
    xx = 5.5 + (i % 3) * 2.42
    yy = 3.62 + (i // 3) * 0.56
    card(s, xx, yy, 2.28, 0.46, fill=WHITE, line=TEAL_LT)
    txt(s, xx + 0.12, yy + 0.12, 2.1, 0.24, l, size=10.5, bold=True, color=TEAL)

rect(s, 0.62, 5.02, 12.1, 0.02, fill=LINE)
txt(s, 0.62, 5.16, 12.1, 0.28, "AZURE HEALTH DATA SERVICES  \u00b7  physical boundary", size=10.5, bold=True, color=GREEN)
card(s, 0.62, 5.46, 5.95, 1.14, fill=WHITE, line=GREEN)
txt(s, 0.92, 5.62, 5.4, 0.26, "fhir-payera", size=13, bold=True, color=INK, font=MONO)
txt(s, 0.92, 5.9, 5.4, 0.5, "Contoso only. Contracts CT-3456 and CT-7788\nseparated logically by meta.tag.", size=11, color=GREY, line_spacing=1.15)
card(s, 6.78, 5.46, 5.94, 1.14, fill=WHITE, line=GREEN)
txt(s, 7.08, 5.62, 5.4, 0.26, "fhir-payerb", size=13, bold=True, color=INK, font=MONO)
txt(s, 7.08, 5.9, 5.4, 0.5, "Fabrikam only. Separate data store.\nNo route exists from Contoso to here.", size=11, color=GREY, line_spacing=1.15)
notes(s, """
Say: "This is the whole architecture on one slide. Read it top to bottom.

At the top, two payers. Each has its own Entra application registration. They get a token
the SMART Backend Services way - client credentials, no user present, because this is
server to server. That is the standard CMS-0057-F expects.

In the middle, the gateway. One APIM instance, and it is the only path in. Six layers of
policy run on every call and I will walk all six in part six.

At the bottom, two FHIR services. This is the physical boundary. Contoso's data is in one,
Fabrikam's is in the other. Inside Contoso's service, its two contracts are separated
logically by a tag.

The sentence I want you to hold on to is on the bottom right: no route exists from Contoso
to here. Not 'no policy allows it' - no route exists. I will prove that in part five."

If someone asks about BasicV2: "That is a POC tier - it does not support VNet injection.
For production you would be on Premium or Standard V2, and I have that in part seven."
""")

# ============================================================  08 the decision
s = new()
head(s, "3 \u00b7 APIM in the CMS architecture", "Physical where it matters, logical where it does not")

card(s, 0.62, 1.62, 5.95, 2.5, fill=WHITE, line=GREEN)
rect(s, 0.62, 1.62, 5.95, 0.055, fill=GREEN)
txt(s, 0.92, 1.86, 5.4, 0.3, "BETWEEN PAYERS  \u2014  PHYSICAL", size=11.5, bold=True, color=GREEN)
txt(s, 0.92, 2.22, 5.4, 0.9,
    "Separate FHIR service. Separate data store. Separate RBAC.",
    size=15, bold=True, color=INK, line_spacing=1.2)
txt(s, 0.92, 2.86, 5.4, 1.1,
    "A defect in an APIM policy cannot leak Contoso to Fabrikam, because Fabrikam\u2019s "
    "credential has no path to Contoso\u2019s service at all. The blast radius of a mistake "
    "is bounded by architecture, not by correctness.",
    size=12.5, color=GREY, line_spacing=1.25)

card(s, 6.78, 1.62, 5.94, 2.5, fill=WHITE, line=AMBER)
rect(s, 6.78, 1.62, 5.94, 0.055, fill=AMBER)
txt(s, 7.08, 1.86, 5.4, 0.3, "WITHIN A PAYER  \u2014  LOGICAL", size=11.5, bold=True, color=AMBER)
txt(s, 7.08, 2.22, 5.4, 0.9,
    "One service. Contracts as meta.tag, enforced at the gateway.",
    size=15, bold=True, color=INK, line_spacing=1.2)
txt(s, 7.08, 2.86, 5.4, 1.1,
    "A defect here could expose one contract to another inside the same payer. "
    "That is a real risk and we accept it deliberately \u2014 the counterparty already "
    "holds a contract with you, and the data is theirs to hold.",
    size=12.5, color=GREY, line_spacing=1.25)

card(s, 0.62, 4.28, 12.1, 0.86, fill=NAVY, line=NAVY)
txt(s, 0.92, 4.5, 11.5, 0.44,
    "Physical separation stops a mistake.  Logical separation stops a query.",
    size=19, bold=True, color=WHITE, font=HEAD, align=PP_ALIGN.CENTER)

txt(s, 0.62, 5.35, 12.1, 0.3, "WHAT THIS BUYS YOU", size=11, bold=True, color=TEAL)
for i, (a, b) in enumerate([
        ("Scales to 200 payers", "You do not provision 200 services. You provision a manageable number and partition contracts inside them."),
        ("Puts cost where risk is", "The expensive boundary is spent only where a failure would cross a company line.")]):
    card(s, 0.62 + i * 6.16, 5.68, 5.95, 0.92, fill=WHITE)
    txt(s, 0.92 + i * 6.16, 5.84, 5.4, 0.26, a, size=13, bold=True, color=INK)
    txt(s, 0.92 + i * 6.16, 6.12, 5.5, 0.44, b, size=11.5, color=GREY, line_spacing=1.15)
notes(s, """
This is the conceptual centre of the deck. Slow down.

Say: "This is the decision this room made on the twelfth, and I want to make the reasoning
explicit because it is the thing you will have to defend to your security review.

Between payers, the boundary is physical. Separate FHIR service, separate data store,
separate access control. The important property is on the left in grey: a defect in one of
my policies cannot leak Contoso to Fabrikam, because Fabrikam's credential has no path to
Contoso's service at all. The blast radius is bounded by architecture rather than by
whether I wrote the policy correctly.

Within a payer, the boundary is logical. Contracts are a tag on the data, enforced here at
the gateway. And I want to be straight with you - a defect here could expose one contract
to another inside the same payer. That is a genuine risk and we are accepting it on
purpose."

Then the navy bar, slowly: "Physical separation stops a mistake. Logical separation stops
a query."

"So you spend the expensive boundary only where a failure crosses a company line, and you
use the cheap one where both parties already have a contract with you. That is what lets
this scale to two hundred payers without provisioning two hundred services."

Expected challenge - "why is logical good enough for contracts?" Answer: "Because the
counterparty is the same legal entity either way. If CT-3456 leaked to CT-7788, both are
Contoso. If Contoso leaked to Fabrikam, that is a breach between two competitors. Different
consequence, different control."
""")

# ================================================  09 inbound vs outbound model
s = new()
head(s, "4 \u00b7 Inbound vs outbound partitioning", "One service, two doors, two different locks")
txt(s, 0.62, 1.5, 12.1, 0.3,
    "Platform Engineering's question: can inbound and outbound be isolated within one service? "
    "Yes \u2014 by not letting either one be reached directly.",
    size=13, color=GREY)

txt(s, 0.62, 1.98, 5.95, 0.28, "OUTBOUND  \u00b7  payer reads", size=11.5, bold=True, color=TEAL)
card(s, 0.62, 2.28, 5.95, 2.5, fill=WHITE, line=TEAL_LT)
txt(s, 0.92, 2.44, 5.4, 0.26, "/payera/outbound   /payerb/outbound", size=12, bold=True, color=INK, font=MONO)
for i, (t, ok) in enumerate([("GET search \u2014 _tag forced to own contracts", True),
                             ("GET Group/{id}/$export \u2014 own Groups only", True),
                             ("GET metadata, smart-configuration", True),
                             ("POST / PUT / PATCH / DELETE", False),
                             ("System or Patient level $export", False)]):
    yy = 2.78 + i * 0.38
    txt(s, 0.95, yy, 0.3, 0.26, "\u2713" if ok else "\u2715", size=13, bold=True,
        color=GREEN if ok else RED)
    txt(s, 1.28, yy, 5.1, 0.26, t, size=11.5, color=INK if ok else GREY)

txt(s, 6.78, 1.98, 5.94, 0.28, "INBOUND  \u00b7  Northwind Health writes", size=11.5, bold=True, color=AMBER)
card(s, 6.78, 2.28, 5.94, 2.5, fill=WHITE, line=AMBER)
txt(s, 7.08, 2.44, 5.4, 0.26, "/payera/inbound   /payerb/inbound", size=12, bold=True, color=INK, font=MONO)
for i, (t, ok) in enumerate([("POST / PUT \u2014 stamps meta.tag on write", True),
                             ("X-Payer-Contract header mandatory", True),
                             ("Bundle transaction ingest", True),
                             ("$export \u2014 denied outright", False),
                             ("Payer credentials \u2014 not entitled here", False)]):
    yy = 2.78 + i * 0.38
    txt(s, 7.11, yy, 0.3, 0.26, "\u2713" if ok else "\u2715", size=13, bold=True,
        color=GREEN if ok else RED)
    txt(s, 7.44, yy, 5.1, 0.26, t, size=11.5, color=INK if ok else GREY)

card(s, 0.62, 4.95, 12.1, 1.65, fill=RGBColor(0xEC, 0xF4, 0xF8), line=TEAL)
txt(s, 0.92, 5.14, 11.5, 0.3, "THE MECHANISM \u2014 why this cannot be worked around", size=11.5, bold=True, color=TEAL)
txt(s, 0.92, 5.5, 11.6, 0.98,
    "Both routes point at the same FHIR service. Neither is reachable directly, because the payer\u2019s "
    "credential holds no role on that service. The route is not a suggestion about which door to use \u2014 "
    "it is the only door that exists, and each door runs a different allow-list before it opens.",
    size=13.5, color=INK, line_spacing=1.28)
notes(s, """
Say: "This is the direct answer to Platform Engineering's question, so I want to be precise.

The question was: can inbound and outbound be isolated from each other within one service?
The answer is yes, and the reason is slightly counter-intuitive - it is because neither of
them can be reached directly.

On the left, the outbound route. A payer can search, and their search is always filtered to
their own contracts whether they ask for it or not. They can run a Group-scoped bulk export,
but only against a Group in their own contract. They can read the capability statement.
They cannot write anything - every write verb is refused. And they cannot run a system-level
or patient-level export, only Group-scoped.

On the right, the inbound route, which is how Northwind Health pushes data in. It can write, and
every write is stamped with a contract tag. That header is mandatory - if it is missing, the
write is rejected, because if ingest does not stamp then export cannot filter. And export is
denied outright on this route, so even if an ingest credential were compromised, it cannot
be turned around and used to read data out."

Then the blue box, which is the part that actually answers the question:

"Both of these point at the same FHIR service. The reason they stay separate is that neither
is reachable directly - the payer's credential holds no role on the service at all. So the
route is not a suggestion about which door to use. It is the only door that exists, and each
door runs a different allow-list before it opens."

Expected question - "what if someone gets the FHIR URL?" Answer: "They can have it. It is
not a secret. Without a role assignment the URL is useless, and that is assertion twelve."
""")

# =================================================  10 how contracts are enforced
s = new()
head(s, "4 \u00b7 Inbound vs outbound partitioning", "How a contract boundary is actually enforced")

txt(s, 0.62, 1.52, 12.1, 0.28, "1  \u00b7  ON INGEST \u2014 every resource is stamped", size=11.5, bold=True, color=AMBER)
code(s, 0.62, 1.84, 5.95, 1.5, [
    ("POST /payera/inbound/Patient", TEAL_LT),
    ("X-Payer-Contract: CT-3456", ICE),
    "",
    ("policy adds \u2192", MUTED),
    ('"meta": { "tag": [{ "system": ".../contract",', ICE),
    ('              "code": "CT-3456" }] }', ICE),
], size=10.5)

txt(s, 6.78, 1.52, 5.94, 0.28, "2  \u00b7  ON READ \u2014 the filter is forced, not requested", size=11.5, bold=True, color=TEAL)
code(s, 6.78, 1.84, 5.94, 1.5, [
    ("GET /payera/outbound/Patient?_tag=ANYTHING", TEAL_LT),
    "",
    ("policy overwrites \u2192", MUTED),
    ("?_tag=.../contract|CT-3456,.../contract|CT-7788", ICE),
    "",
    ("caller-supplied _tag is discarded, never merged", RGBColor(0xF0, 0xC8, 0x7A)),
], size=10.5)

card(s, 0.62, 3.56, 12.1, 1.28, fill=WHITE, line=LINE)
txt(s, 0.92, 3.74, 11.5, 0.28, "THE ENTITLEMENT MAP  \u00b7  an APIM named value, editable without redeploying", size=11.5, bold=True, color=TEAL)
txt(s, 0.92, 4.06, 11.6, 0.66,
    "cmsdqm-payera  \u2192  payer 'payera', contracts CT-3456 + CT-7788, groups group-ct3456 + group-ct7788\n"
    "cmsdqm-payerb  \u2192  payer 'payerb', contract CT-9001, group group-ct9001",
    size=12, color=INK, font=MONO, line_spacing=1.35)

card(s, 0.62, 5.0, 5.95, 1.6, fill=GREEN_LT, line=GREEN)
txt(s, 0.92, 5.2, 5.4, 0.28, "ONBOARDING PAYER #3", size=11.5, bold=True, color=GREEN)
txt(s, 0.92, 5.54, 5.5, 0.9,
    "Register an app, add one line to the entitlement map, create the Group. "
    "No redeployment, no code change, no downtime.",
    size=12.5, color=INK, line_spacing=1.25)

card(s, 6.78, 5.0, 5.94, 1.6, fill=WHITE, line=LINE)
txt(s, 7.08, 5.2, 5.4, 0.28, "WHY A TAG AND NOT A SEARCH FILTER", size=11.5, bold=True, color=TEAL)
txt(s, 7.08, 5.54, 5.5, 0.9,
    "The tag travels with the resource. Any future consumer \u2014 export, analytics, a "
    "second gateway \u2014 sees the same boundary without re-deriving it.",
    size=12.5, color=GREY, line_spacing=1.25)
notes(s, """
Say: "Two steps, and they only work as a pair.

On ingest, every resource gets stamped. Northwind Health posts a Patient with a contract header,
and the policy writes that contract into meta.tag on the resource itself. That header is
mandatory - I mentioned that on the last slide, and this is why. If ingest does not stamp,
export has nothing to filter on.

On read, the filter is forced. Look carefully at the top line on the right - the payer sent
a _tag of their own, and it does not matter what they put there. The policy overwrites it.
It does not merge it, it does not validate it, it discards it and substitutes the contracts
this caller actually holds. A payer cannot widen their own scope by editing the query
string. That is assertion thirteen and I will show it running.

In the middle is the entitlement map. That is a named value in APIM - configuration, not
code. You are looking at the real contents, minus the secrets.

And that gives you the green box, which is the operational answer to 'what does payer number
three cost us'. Register an app, add one line to that map, create the Group. No
redeployment, no code change, no downtime."

Expected question - "why a tag rather than just filtering in the query?" Point at the box
on the right: "Because the tag travels with the resource. Any future consumer - a bulk
export, an analytics pipeline, a second gateway - sees the same boundary without having to
re-derive it. The boundary is in the data, not just in the request path."
""")

# ==================================================  11 isolation tests overview
s = new()
head(s, "5 \u00b7 Gateway isolation with tests", "Sixteen assertions \u2014 the guarantees, written as code")
txt(s, 0.62, 1.5, 12.1, 0.3,
    "Each test names the status code it expects. A green run is the argument; the documentation is only commentary on it.",
    size=13, color=GREY)

groups = [
    ("BASELINE \u2014 the happy path, or the negatives prove nothing", GREEN, [
        ("1", "own data readable", "200"), ("2", "Group export accepted", "202"),
        ("3", "capability statement", "200")]),
    ("CROSS-PAYER \u2014 the boundary that must never bend", RED, [
        ("4", "payer B app, valid payer A audience", "403"),
        ("5", "payer B token, payer B audience", "401"),
        ("12", "payer token straight to AHDS", "403")]),
    ("ROUTE + SCOPE \u2014 inbound / outbound and contract limits", TEAL, [
        ("6", "unentitled Group export", "403"), ("7", "write on outbound route", "403"),
        ("8", "payer credential on inbound route", "403"), ("9", "export on inbound route", "403"),
        ("10", "system-level export", "403"), ("11", "patient-level export", "403")]),
    ("DATA + LOAD \u2014 we check the body, not just the code", AMBER, [
        ("13", "caller _tag overridden", "200"), ("13b", "body has only own contracts", "CT-3456"),
        ("14", "untagged inbound write rejected", "403"), ("15", "second export in 5 min", "429")]),
]
y = 1.94
for title, col, tests in groups:
    h = 0.42 + 0.3 * ((len(tests) + 2) // 3)
    card(s, 0.62, y, 12.1, h, fill=WHITE)
    rect(s, 0.62, y, 0.055, h, fill=col)
    txt(s, 0.92, y + 0.11, 11.6, 0.26, title, size=11.5, bold=True, color=col)
    for i, (num, name, exp) in enumerate(tests):
        xx = 0.95 + (i % 3) * 3.95
        yy = y + 0.42 + (i // 3) * 0.3
        txt(s, xx, yy, 0.5, 0.24, num, size=11, bold=True, color=GREY, font=MONO)
        txt(s, xx + 0.42, yy, 2.6, 0.24, name, size=11, color=INK)
        txt(s, xx + 3.02, yy, 0.7, 0.24, exp, size=11, bold=True, color=col, font=MONO)
    y += h + 0.14
card(s, 0.62, 5.86, 12.1, 0.78, fill=RGBColor(0xEC, 0xF4, 0xF8), line=TEAL)
txt(s, 0.92, 6.04, 11.6, 0.46,
    [("./scripts/run-isolation-tests.ps1", {"bold": True, "size": 13, "color": INK, "font": MONO}),
     ("     runs against the deployed gateway  \u00b7  mints short-lived credentials and destroys them on exit",
      {"size": 12.5, "color": GREY})], size=12.5)
notes(s, """
This slide is the set-up for the live run. Do not linger - ninety seconds, then switch to
the terminal.

Say: "Before I run this, here is what is in it, grouped four ways.

The baseline group matters more than it looks. If the happy path does not work, none of the
negative tests prove anything - a gateway that refuses everything would pass all my security
assertions and be useless. So first we prove a payer can read their own data and start a
bulk export.

The cross-payer group is the boundary that must never bend. Three tests, and I will stop on
two of them when we run it.

The route and scope group is Platform Engineering's question, made executable. Writes refused on
the outbound route, exports refused on the inbound route, system and patient level exports
refused everywhere.

And the last group is the one I am proudest of, because thirteen-b does not check a status
code - it reads the response body and asserts that only CT-3456 came back. A two hundred
would have passed. We check the data."

Then: "Let me just run it." Switch to the terminal.
""")

# =================================================  12 the three that matter
s = new(dark=True)
head(s, "5 \u00b7 Gateway isolation with tests", "The three lines to stop on", dark=True)
txt(s, 0.85, 1.5, 11.8, 0.3, "Sixteen will scroll past. Narrate these and no others.",
    size=13, color=MUTED)

three = [
    ("4", "Valid token. Wrong payer.", "403",
     "Fabrikam's application, holding a technically valid token minted for Contoso's audience. "
     "The signature checks out. The entitlement says payerb, the endpoint says payera, and the call dies at layer 2.",
     "The boundary holds against a legitimate credential, not just a forged one."),
    ("12", "Payer goes straight to AHDS.", "403",
     "Bypassing the gateway entirely \u2014 payer token, FHIR hostname, no APIM in the path. "
     "Entra authenticates the caller perfectly. AHDS then finds no role assignment.",
     "403, not 401. Skipping the gateway does not skip the control \u2014 it removes the only means of access."),
    ("13b", "We read the body, not the code.", "CT-3456",
     "Payer A holds two contracts. The response is inspected resource by resource and asserted to "
     "contain CT-3456 only, because the Group being exported belongs to that contract alone.",
     "A 200 would have passed a status-code check. Logical separation is verified in the data."),
]
y = 1.96
for num, title, code_, detail, proves in three:
    card(s, 0.85, y, 11.65, 1.5, fill=NAVY_CARD, line=SLATE)
    chip(s, 1.08, y + 0.28, 0.44, num, fill=TEAL, size=13)
    txt(s, 1.72, y + 0.2, 6.2, 0.3, title, size=15, bold=True, color=WHITE)
    txt(s, 1.72, y + 0.54, 8.6, 0.56, detail, size=11.5, color=ICE, line_spacing=1.18)
    txt(s, 1.72, y + 1.14, 8.6, 0.28, proves, size=11.5, bold=True, color=TEAL_LT)
    txt(s, 10.6, y + 0.52, 1.7, 0.4, code_, size=17, bold=True, color=GREEN,
        font=MONO, align=PP_ALIGN.RIGHT)
    y += 1.62
notes(s, """
Use this slide only if the terminal output scrolled too fast, or as the summary right after
the run. It is the same three lines.

Say: "Sixteen went past. Three of them carry the argument.

Number four. Fabrikam's application, holding a token that is technically valid - correctly
signed, correct audience for Contoso's endpoint. It still gets a 403. The entitlement map
says this caller belongs to payer B, the endpoint is payer A, and the call dies at layer two.
What that proves is that the boundary holds against a legitimate credential. Anyone can stop
a forged token. Stopping a real one that is simply pointed at the wrong tenant is the harder
case.

Number twelve is the one I would put in front of your security review. Here the payer skips
the gateway completely - payer token, FHIR hostname, no APIM anywhere in the path. And the
answer is 403, not 401.

That distinction is the whole design. A 401 would mean 'we do not recognise you'. A 403
means Entra recognised them perfectly and they simply hold no permission. And no policy bug
of mine can grant it, because there is no policy in that path at all.

Thirteen-b. Payer A holds two contracts. We do not check the status code, we read the body,
resource by resource, and assert that only CT-3456 came back. A 200 would have passed a
status check and told us nothing."

Land on: "Skipping the gateway does not skip the control. It removes the only means of
access."
""")

# ==================================================  13 policy walkthrough
s = new()
head(s, "6 \u00b7 Policy walkthrough", "Six layers, in the order they execute")
txt(s, 0.62, 1.5, 12.1, 0.28,
    "apim/policies/payer-outbound.xml  \u00b7  378 lines  \u00b7  line numbers are live \u2014 open it and follow along",
    size=12, color=GREY, font=MONO)

rows = [
    ("1", "Authentication", "27", "validate-jwt against Entra. Signature, expiry, audience, issuer.", "401 if the token is not ours", TEAL),
    ("2", "Entitlement", "47", "Resolve azp claim \u2192 payer + contracts. Cross-payer guard here.", "403 if unknown or wrong payer", TEAL),
    ("3", "Route allow-list", "124", "Classify the operation. Outbound is read-only, full stop.", "403 on any write verb", AMBER),
    ("4a", "Group scoping", "165", "$export must name a Group inside your own contract.", "403 on system / patient / other", AMBER),
    ("4b", "Tag injection", "239", "_tag overwritten with the caller's contracts. Never merged.", "Scope cannot be widened", AMBER),
    ("5", "Rate + lock", "262", "600/min, 50k/day, and one bulk export per payer per 5 min.", "429 on the second export", RED),
    ("6", "Trusted broker", "317", "Managed identity token attached. Payer token discarded.", "The payer never reaches FHIR", GREEN),
]
y = 1.92
for num, name, ln, what, effect, col in rows:
    card(s, 0.62, y, 12.1, 0.63, fill=WHITE)
    rect(s, 0.62, y, 0.055, 0.63, fill=col)
    txt(s, 0.9, y + 0.19, 0.5, 0.26, num, size=13, bold=True, color=col, font=MONO)
    txt(s, 1.45, y + 0.19, 2.0, 0.26, name, size=12.5, bold=True, color=INK)
    txt(s, 3.5, y + 0.2, 0.75, 0.24, f"L{ln}", size=11, color=GREY, font=MONO)
    txt(s, 4.35, y + 0.19, 4.9, 0.26, what, size=11.5, color=INK)
    txt(s, 9.4, y + 0.19, 3.2, 0.26, effect, size=11, bold=True, color=col, align=PP_ALIGN.RIGHT)
    y += 0.7
notes(s, """
Have the policy XML open in the portal on the second screen for this. Scroll to each line
number as you say it. Do not read all six in depth - two minutes total, then spend the time
on layer 6.

Say: "This is the outbound policy. Three hundred and seventy-eight lines, and the line
numbers on screen are real - you can open this file and follow along.

Layer one authenticates. Standard JWT validation against Entra.

Layer two is entitlement - it takes the azp claim, which is the calling application's ID,
and looks up which payer and which contracts that application holds. The cross-payer guard
lives here, and that is what assertion four hits.

Layer three classifies the operation and enforces the allow-list. Outbound is read-only.

Four-a and four-b are the contract boundary - Group scoping for exports, tag injection for
searches.

Layer five is rate limiting, and there is a detail there worth knowing: the export lock is
a cache entry rather than a second rate-limit policy, because two rate-limit policies in one
section do not both fire, and the failure is silent. That is a real APIM behaviour and it
cost me an afternoon."

Then slow down for layer 6:

"And layer six is the keystone. The gateway attaches its own managed identity token and
discards the payer's. Every call that reaches AHDS is made by the gateway, as the gateway.
The payer's token got them through the front door and no further."

If asked why 401 at layer 1 but 403 at layer 2: "401 means we do not know who you are. 403
means we know exactly who you are and you may not do this. Getting that right matters for
the payer's own error handling."
""")

# =====================================================  14 debug trace
s = new()
head(s, "6 \u00b7 Debug trace", "Watching one call go through the pipeline")

card(s, 0.62, 1.58, 5.95, 2.15, fill=WHITE, line=LINE)
txt(s, 0.92, 1.76, 5.4, 0.28, "HOW TO TURN IT ON", size=11.5, bold=True, color=TEAL)
txt(s, 0.92, 2.1, 5.5, 1.5,
    "Portal \u2192 APIM \u2192 APIs \u2192 the operation \u2192 Test tab \u2192 Trace.\n\n"
    "Or send the call yourself with an Ocp-Apim-Trace header and read the trace "
    "location returned in the response.",
    size=12.5, color=GREY, line_spacing=1.3)

card(s, 6.78, 1.58, 5.94, 2.15, fill=AMBER_LT, line=AMBER)
txt(s, 7.08, 1.76, 5.4, 0.28, "WHY IT MATTERS OPERATIONALLY", size=11.5, bold=True, color=AMBER)
txt(s, 7.08, 2.1, 5.5, 1.5,
    "When a payer says \u201cyour API is broken\u201d, this tells you in seconds whether "
    "you refused them, and at which layer \u2014 without turning on verbose logging or "
    "asking them to reproduce it.",
    size=12.5, color=INK, line_spacing=1.3)

txt(s, 0.62, 3.94, 12.1, 0.28, "WHAT THE TRACE SHOWS \u2014 a refused cross-payer call", size=11.5, bold=True, color=TEAL)
code(s, 0.62, 4.26, 12.1, 2.3, [
    ("inbound", TEAL_LT),
    ("  validate-jwt          token valid, audience matched", ICE),
    ("  set-variable          callerAppId = bbbbbbbb-\u2026  (payer B application)", ICE),
    ("  set-variable          entitlement = { payer: 'payerb', contracts: ['CT-9001'] }", ICE),
    ("  choose \u2192 when         payerKey 'payerb' != 'payera'", RGBColor(0xF0, 0xC8, 0x7A)),
    ("  return-response       403  OperationOutcome", RGBColor(0xE8, 0x94, 0x8A)),
    "",
    ("backend                 not executed \u2014 AHDS was never called", MUTED),
], size=11)
notes(s, """
This is the part of the session that is new to most people, and it lands well. Budget three
minutes.

Say: "Last thing on policy, and this one is operational rather than architectural.

API Management can trace a single call through the pipeline. You turn it on in the portal
under the Test tab, or you send the call yourself with a trace header.

What you get is on the bottom of the slide, and this is a real refused call - Fabrikam's
application against Contoso's endpoint, which is assertion four.

Read down it. The token validated fine. The caller was resolved to Fabrikam's application
ID. Their entitlement came back as payer B with contract CT-9001. Then the cross-payer guard
compared payerb against payera, they did not match, and the call returned a 403 with a FHIR
OperationOutcome.

And look at the last line. The backend section never executed. AHDS was never called. The
refusal happened at the gateway and Contoso's FHIR service never saw a request at all."

Then the operational point, which is the one they will actually care about:

"Why this matters day to day - when a payer opens a ticket saying your API is broken, this
tells you in seconds whether you refused them and at which layer. You do not need verbose
logging on, and you do not need to ask them to reproduce it."

If asked about retention: "Traces are short-lived and only generated on request. For
persistent evidence you want the gateway logs in Log Analytics - there is a KQL query for
403s by payer in the environment script."
""")

# ==============================================  15 credentials and identity
s = new()
head(s, "7 \u00b7 Credentials and identity", "How a payer authenticates, and what you run in production")

txt(s, 0.62, 1.5, 12.1, 0.28, "THE PAYER SIDE  \u00b7  SMART on FHIR Backend Services", size=11.5, bold=True, color=TEAL)
opts = [
    ("Client secret", "POC only", "Simple, but it is a shared password with an expiry you must track. "
     "This tenant caps them at 30 days.", RED),
    ("Certificate credential", "Workable", "Payer holds a private key and signs a client assertion. "
     "No shared secret in flight. Rotation is still yours to run.", AMBER),
    ("Federated identity credential", "Preferred", "Payer's own IdP issues the assertion; Entra trusts it. "
     "No secret and no certificate to rotate at all.", GREEN),
]
y = 1.82
for name, verdict, desc, col in opts:
    card(s, 0.62, y, 12.1, 0.82, fill=WHITE)
    rect(s, 0.62, y, 0.055, 0.82, fill=col)
    txt(s, 0.95, y + 0.14, 3.2, 0.28, name, size=13, bold=True, color=INK)
    txt(s, 0.95, y + 0.44, 3.2, 0.26, verdict, size=11, bold=True, color=col)
    txt(s, 4.3, y + 0.22, 8.2, 0.44, desc, size=12, color=GREY, line_spacing=1.15)
    y += 0.92

txt(s, 0.62, 4.68, 12.1, 0.28, "THE NORTHWIND HEALTH SIDE  \u00b7  what actually holds the FHIR role", size=11.5, bold=True, color=GREEN)
card(s, 0.62, 5.0, 12.1, 1.6, fill=GREEN_LT, line=GREEN)
txt(s, 0.92, 5.2, 11.5, 0.32,
    "Only the APIM managed identity holds FHIR Data Contributor. Verified live:",
    size=13, bold=True, color=INK)
txt(s, 0.92, 5.58, 11.6, 0.92,
    "cmsdqm-payera   aaaaaaaa-\u2026     0 Azure role assignments\n"
    "cmsdqm-payerb   bbbbbbbb-\u2026     0 Azure role assignments\n"
    "APIM managed identity      dddddddd-\u2026     FHIR Data Contributor on both services",
    size=12, color=INK, font=MONO, line_spacing=1.35)
notes(s, """
Say: "Two halves to this. How the payer proves who they are, and what actually holds
permission on our side.

Top half is the payer. All three of these are the SMART Backend Services pattern - client
credentials, no user present. The difference is only what the payer holds.

A client secret is a shared password. It works, it is what the POC uses, and I would not
ship it - partly because this tenant caps secrets at thirty days, which becomes an
operational treadmill across forty payers.

A certificate credential means the payer signs an assertion with a private key. Nothing
shared travels. Better, but you still own a rotation process.

A federated identity credential is where I would land. The payer's own identity provider
issues the assertion and Entra is configured to trust it. There is no secret and no
certificate to rotate, because the payer is already managing that lifecycle for themselves.

Bottom half is the important one, and these numbers are from the live environment this
morning. Both payer applications hold zero Azure role assignments. Zero. The only thing with
FHIR Data Contributor is the gateway's managed identity."

Land it: "That is why assertion twelve returns a 403. There is nothing to take away from
those applications, because they were never given anything."

Expected question - "how do we rotate forty payer secrets?" Answer: "You do not. You move to
federated credentials and the payer manages it. That is decision three on the next slide."
""")

# ==============================================  16 deployment options
s = new()
head(s, "7 \u00b7 Deployment options", "What you would actually run")

txt(s, 0.62, 1.5, 12.1, 0.28, "TIER \u2014 the POC is on BasicV2, which is not a production answer", size=11.5, bold=True, color=TEAL)
tiers = [
    ("BasicV2", "POC only", "No VNet injection. Fine for proving behaviour, not for PHI in production.", RED),
    ("StandardV2", "Likely fit", "VNet integration, autoscale. Sufficient unless you need full network isolation.", GREEN),
    ("Premium", "If required", "VNet injection, multi-region, availability zones. Choose if security review demands it.", AMBER),
]
y = 1.82
for name, verdict, desc, col in tiers:
    card(s, 0.62, y, 12.1, 0.72, fill=WHITE)
    rect(s, 0.62, y, 0.055, 0.72, fill=col)
    txt(s, 0.95, y + 0.11, 2.6, 0.28, name, size=13, bold=True, color=INK)
    txt(s, 0.95, y + 0.4, 2.6, 0.26, verdict, size=11, bold=True, color=col)
    txt(s, 3.9, y + 0.22, 8.6, 0.3, desc, size=12, color=GREY)
    y += 0.82

txt(s, 0.62, 4.38, 12.1, 0.28, "THE UNQUANTIFIED RISK \u2014 and the one thing I need from you", size=11.5, bold=True, color=AMBER)
card(s, 0.62, 4.7, 12.1, 1.9, fill=AMBER_LT, line=AMBER)
txt(s, 0.92, 4.92, 11.5, 0.34,
    "\u201cUsually happens at the same time \u2026 will pop the server itself.\u201d  \u2014  Enterprise Architecture, 8/12",
    size=13, italic=True, color=INK)
txt(s, 0.92, 5.34, 11.6, 1.1,
    "Forty payers exporting on the same overnight schedule is the failure mode this design has not yet been "
    "sized against. The gateway holds one export per payer per five minutes, which converts a stampede into a "
    "queue \u2014 but the right window, and the right AHDS capacity behind it, needs your real numbers: "
    "how many payers, how many members each, and on what schedule.",
    size=13, color=INK, line_spacing=1.28)
notes(s, """
Say: "Two things here. What tier, and the risk I cannot close on my own.

On tier - the POC runs on BasicV2 and I want to be clear that is not a production answer. It
does not support VNet injection. It was the right call for proving behaviour quickly and the
wrong call for PHI in production.

StandardV2 is probably your fit - VNet integration and autoscale. Premium if your security
review insists on full VNet injection, multi-region or availability zones. That is a
conversation with your security architects, not a technical blocker either way. The policies
are identical across all three - nothing I have shown you today changes.

Now the amber box, and this is the most important slide in part seven."

Read the quote. Then:

"This is the failure mode this design has not been sized against. Forty payers all exporting
on the same overnight schedule. The gateway holds one export per payer per five minutes,
which turns a stampede into a queue - that is assertion fifteen, the 429. But I picked five
minutes as a reasonable-sounding number, not as the output of a capacity model.

To size it properly I need three numbers from you: how many payers, how many members each,
and on what schedule they export. Give me those and I will come back with a sized design
and a cost. Without them, anything I tell you about capacity is a guess with a confident
tone."

That honesty plays well. Do not oversell here.
""")

# ==============================================  17 decisions
s = new()
head(s, "8 \u00b7 Decisions", "Four things I need from this room")
decisions = [
    ("1", "Confirm the separation model",
     "Physical per payer, logical per contract, enforced at the gateway.",
     "Everything downstream assumes it. If this is wrong, better to know now.", TEAL),
    ("2", "Confirm inbound and outbound stay separate APIs",
     "Two routes over one FHIR service, with mirrored allow-lists.",
     "The alternative is two FHIR services per payer, which doubles the estate.", TEAL),
    ("3", "Choose the payer credential model",
     "Client secret, certificate, or federated identity credential.",
     "This one has a long lead time \u2014 it changes payer onboarding paperwork.", AMBER),
    ("4", "Give me the export volume numbers",
     "Payers, members per payer, and export schedule.",
     "Blocks capacity sizing and therefore blocks the production cost estimate.", RED),
]
y = 1.7
for num, title, detail, why, col in decisions:
    card(s, 0.62, y, 12.1, 1.16, fill=WHITE)
    rect(s, 0.62, y, 0.055, 1.16, fill=col)
    chip(s, 0.92, y + 0.34, 0.46, num, fill=col, size=14)
    txt(s, 1.62, y + 0.16, 8.0, 0.3, title, size=14.5, bold=True, color=INK)
    txt(s, 1.62, y + 0.5, 10.6, 0.28, detail, size=12.5, color=INK)
    txt(s, 1.62, y + 0.8, 10.6, 0.28, why, size=11.5, italic=True, color=GREY)
    y += 1.28
notes(s, """
Slow down. This is what the session was for.

Say: "Four decisions, and I would like to leave with at least the first two.

One - confirm the separation model. Physical per payer, logical per contract, enforced at
the gateway. Everything downstream assumes this, so if it is wrong I would much rather find
out in this room than after we have built on it.

Two - confirm inbound and outbound stay as separate APIs over one FHIR service. The
alternative is two FHIR services per payer, which doubles your estate and your cost. I think
what I showed you in part four makes the case, but it is your call.

Three - the payer credential model. This one has a long lead time because it changes the
paperwork you send a payer during onboarding. My recommendation is federated identity
credentials, but I understand if you want to start with certificates.

Four - and this is the one that actually blocks me. I need your export volume numbers.
Payers, members per payer, schedule. Until I have those I cannot size production and I
cannot give you a real cost."

Then stop talking and let them respond. Do not fill the silence.

If they will not decide today: "That is fine - can we put a date on each? I will take an
owner and a date rather than a decision."
""")

# ==============================================  18 next steps
s = new(dark=True)
head(s, "8 \u00b7 Next steps", "What happens after this call", dark=True)

txt(s, 0.85, 1.56, 11.8, 0.28, "IMMEDIATE  \u00b7  this week", size=11, bold=True, color=TEAL_LT)
for i, (who, what) in enumerate([
        ("Northwind Health", "Export volume numbers \u2014 payers, members, schedule"),
        ("Northwind Health", "Security review of the separation model as shown"),
        ("Microsoft", "Sized production design and cost once volumes land"),
        ("Microsoft", "Federated credential onboarding guide for payers")]):
    yy = 1.88 + i * 0.56
    card(s, 0.85, yy, 11.65, 0.48, fill=NAVY_CARD, line=SLATE)
    txt(s, 1.08, yy + 0.12, 1.9, 0.26, who, size=11.5, bold=True,
        color=TEAL_LT if who == "Microsoft" else RGBColor(0xF0, 0xC8, 0x7A))
    txt(s, 3.1, yy + 0.12, 9.2, 0.26, what, size=12.5, color=WHITE)

txt(s, 0.85, 4.3, 11.8, 0.28, "THE ENVIRONMENT STAYS UP", size=11, bold=True, color=TEAL_LT)
card(s, 0.85, 4.62, 11.65, 1.0, fill=NAVY_CARD, line=SLATE)
txt(s, 1.08, 4.82, 11.2, 0.62,
    "rg-ahds-fhir-poc remains deployed and is yours to poke at. The policies, the Bicep and the test suite "
    "are all in the handover folder \u2014 run the suite yourself, change a policy, watch a test go red.",
    size=13, color=ICE, line_spacing=1.25)

txt(s, 0.85, 5.86, 11.65, 0.5,
    "The strongest thing you can do with this is break it. If you can get one payer\u2019s data out of the other\u2019s endpoint, I want to hear about it.",
    size=14.5, bold=True, color=WHITE, font=HEAD, align=PP_ALIGN.CENTER)
foot(s, dark=True)
notes(s, """
Close here. Keep it short - people are ready to leave.

Say: "Four actions, two on each side.

You owe me export volumes and a security review. I owe you a sized production design once
those volumes land, and an onboarding guide for federated credentials.

The environment stays up. The resource group, the policies, the Bicep and the test suite are
all in the handover folder. Run the suite yourself. Change a policy and watch a test go red -
that is genuinely the fastest way to understand what the gateway is doing.

And the last line is a real invitation, not a rhetorical one. The strongest thing you can do
with this is try to break it. If anyone here can get one payer's data out of the other's
endpoint, I want to hear about it before your security review does."

Then: "What did I miss?" And stop.
""")

out = Path(__file__).with_name("Prov-APIM-FHIR-Session-2026-09-01.pptx")
prs.save(out)

missing = [i + 1 for i, sl in enumerate(prs.slides)
           if not sl.notes_slide.notes_text_frame.text.strip()]
words = sum(len(sl.notes_slide.notes_text_frame.text.split()) for sl in prs.slides)
print(f"Deck written: {out}")
print(f"Slides: {len(prs.slides.__iter__.__self__._sldIdLst)}  "
      f"Notes missing on: {missing or 'none'}  Notes words: {words}")
