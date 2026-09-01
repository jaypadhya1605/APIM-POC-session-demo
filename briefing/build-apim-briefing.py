"""Builds the Northwind Health APIM briefing deck, with full speaker notes on every slide.

    %LOCALAPPDATA%\\venvs\\pptxbuild\\Scripts\\python.exe build-apim-briefing.py

Style deliberately matches slides/build-deck.py so the two decks can be shown back to back.
Facts are pulled from the live environment in rg-ahds-fhir-poc (see FACTS below).
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

W, H = 13.333, 7.5

NAVY = RGBColor(0x0E, 0x2A, 0x47)
NAVY_CARD = RGBColor(0x1A, 0x3E, 0x60)
TEAL = RGBColor(0x1C, 0x72, 0x93)
TEAL_LT = RGBColor(0x4F, 0xA6, 0xC4)
ICE = RGBColor(0xCF, 0xE3, 0xEF)
PAPER = RGBColor(0xF5, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x12, 0x28, 0x3D)
GREY = RGBColor(0x5A, 0x6B, 0x7B)
LINE = RGBColor(0xD8, 0xE2, 0xEA)
AMBER = RGBColor(0xD9, 0x82, 0x2B)
AMBER_LT = RGBColor(0xFD, 0xF3, 0xE3)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
GREEN_LT = RGBColor(0xE8, 0xF4, 0xEC)
RED = RGBColor(0xB0, 0x3A, 0x2B)
RED_LT = RGBColor(0xFA, 0xEC, 0xEA)
SLATE = RGBColor(0x2C, 0x50, 0x74)
MUTED = RGBColor(0x9F, 0xB8, 0xCC)

HEAD = "Cambria"
BODY = "Calibri"
MONO = "Consolas"

APIM = "apim-poc-ahds-demo01"
GATEWAY = "https://apim-poc-ahds-demo01.azure-api.net"

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

_n = {"i": 0}


def new(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY if dark else PAPER
    _n["i"] += 1
    return s


def rect(s, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE, radius=None, line=None, lw=1.0):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def txt(s, x, y, w, h, runs, size=14, color=INK, font=BODY, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, italic=False,
        line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    if isinstance(runs, str):
        runs = [(runs, {})]
    runs = [(r, {}) if isinstance(r, str) else r for r in runs]

    for i, (text, ov) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("line_spacing", line_spacing)
        if i > 0:
            p.space_before = Pt(ov.get("space_before", space))
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic)
        f.color.rgb = ov.get("color", color)
    return tb


def foot(s, dark=False):
    rect(s, 0.62, 6.94, 12.1, 0.014, fill=SLATE if dark else LINE)
    txt(s, 0.62, 7.03, 9.6, 0.24,
        "Northwind Health  ·  Azure API Management briefing  ·  AHDS payer data exchange",
        size=9, color=RGBColor(0x7F, 0x9A, 0xB5) if dark else GREY)
    txt(s, 11.0, 7.03, 1.72, 0.24, f"{_n['i']:02d}", size=9, align=PP_ALIGN.RIGHT,
        color=RGBColor(0x7F, 0x9A, 0xB5) if dark else GREY)


def head(s, kicker, title, dark=False):
    rect(s, 0.62, 0.46, 0.075, 0.72, fill=TEAL_LT if dark else TEAL)
    txt(s, 0.85, 0.46, 11.8, 0.26, kicker.upper(), size=10.5, bold=True,
        color=TEAL_LT if dark else TEAL, font=BODY)
    txt(s, 0.85, 0.72, 11.8, 0.52, title, size=27, bold=True,
        color=WHITE if dark else INK, font=HEAD)
    foot(s, dark)


def card(s, x, y, w, h, fill=WHITE, line=LINE, radius=0.06):
    return rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                radius=radius, line=line, lw=0.75)


def chip(s, x, y, d, label, fill=TEAL, color=WHITE, size=12):
    rect(s, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)
    txt(s, x, y + d / 2 - 0.115, d, 0.25, label, size=size, bold=True,
        color=color, align=PP_ALIGN.CENTER, font=BODY)


def code(s, x, y, w, h, lines, size=11.5, fill=RGBColor(0x0B, 0x22, 0x3A)):
    rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    runs = []
    for i, ln in enumerate(lines):
        if isinstance(ln, tuple):
            text, col = ln
        else:
            text, col = ln, ICE
        runs.append((text, {"color": col, "space_before": 0 if i == 0 else 3}))
    txt(s, x + 0.24, y + 0.2, w - 0.48, h - 0.4, runs, size=size, font=MONO,
        line_spacing=1.18)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


# ------------------------------------------------------------------ 01 title
s = new(dark=True)
rect(s, 0, 0, 0.28, H, fill=TEAL)
rect(s, 7.9, 0, 5.44, H, fill=RGBColor(0x0B, 0x22, 0x3A))

rect(s, 8.6, 1.5, 3.9, 1.28, fill=RGBColor(0x14, 0x33, 0x54),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
txt(s, 8.85, 1.72, 3.4, 0.3, "PAYER  →  APIM", size=11, bold=True, color=TEAL_LT, font=MONO)
txt(s, 8.85, 2.08, 3.4, 0.5, "validate · entitle · scope\nthrottle · broker", size=10,
    color=MUTED, font=MONO, line_spacing=1.25)

rect(s, 8.6, 3.02, 1.85, 1.15, fill=RGBColor(0x14, 0x33, 0x54),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
rect(s, 8.6, 3.02, 1.85, 0.1, fill=GREEN)
txt(s, 8.8, 3.28, 1.5, 0.6, "OUTBOUND\nread only", size=10, bold=True, color=ICE,
    font=MONO, line_spacing=1.25)

rect(s, 10.65, 3.02, 1.85, 1.15, fill=RGBColor(0x14, 0x33, 0x54),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
rect(s, 10.65, 3.02, 1.85, 0.1, fill=AMBER)
txt(s, 10.85, 3.28, 1.5, 0.6, "INBOUND\nwrite only", size=10, bold=True, color=ICE,
    font=MONO, line_spacing=1.25)

rect(s, 8.6, 4.42, 3.9, 0.92, fill=RGBColor(0x14, 0x33, 0x54),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
txt(s, 8.85, 4.66, 3.4, 0.5, "ONE FHIR SERVICE\nper payer", size=10.5, bold=True,
    color=TEAL_LT, font=MONO, line_spacing=1.25)

txt(s, 8.6, 5.62, 4.3, 0.7,
    [("Two routes. One data plane.", {"bold": True, "color": ICE}),
     ("The route is the boundary.", {"color": MUTED, "space_before": 4})], size=10.5)

txt(s, 1.0, 1.86, 6.6, 0.3, "AZURE API MANAGEMENT", size=11.5, bold=True, color=TEAL_LT)
txt(s, 1.0, 2.28, 6.6, 1.9,
    [("Inbound and outbound,", {}), ("isolated at the gateway.", {"color": ICE})],
    size=42, bold=True, color=WHITE, font=HEAD, line_spacing=1.02)
rect(s, 1.0, 4.3, 1.5, 0.035, fill=TEAL)
txt(s, 1.0, 4.64, 6.4, 1.1,
    "The briefing requested on 18 August: how one FHIR service per payer serves "
    "payer traffic and ingest traffic without either being able to reach the other — "
    "shown live against a deployed environment, not described.",
    size=13.5, color=RGBColor(0xB9, 0xCD, 0xDD), line_spacing=1.28)
txt(s, 1.0, 6.06, 6.4, 0.3,
    "Microsoft Health & Life Sciences  ·  60 minutes  ·  live demo included",
    size=10.5, bold=True, color=TEAL_LT)
notes(s, """
Open here. Do not start the deck until the demo terminal is already warm in another window.

SAY:
"Thanks for the time. This session exists because of one question Ashay raised on the
eighteenth — inbound versus outbound partitioning — and a commitment Joe made on the same
call to bring an APIM specialist back to walk through it properly.

So that's the deal for the next hour. I'm not going to describe an architecture to you. There
is a deployed environment in the non-production subscription, and about twenty minutes from
now I'm going to run a test suite against it live and you'll see sixteen assertions go green
on screen. If the model is wrong, that goes red in front of you.

One honest scoping note before we start, and Steve made this point on the eighteenth: you do
not need API Management to do plan-level filtering. That works directly against the FHIR
endpoints today. APIM is what you add when you need contract-level scoping, inbound/outbound
separation, and burst control. It's additive. I'd rather say that up front than have you
discover it in month three."
""")

# ------------------------------------------------------------------ 02 agenda
s = new()
head(s, "Run of show", "Sixty minutes, eight topics, one live demo")

rows = [
    ("00–05", "Objectives and scope", "Why we're here; what APIM is and is not for"),
    ("05–12", "APIM in the CMS-0057-F architecture", "Where the gateway sits; the 8/12 decision"),
    ("12–20", "Inbound vs outbound partitioning", "The question from 8/18, answered"),
    ("20–35", "Live demo — gateway isolation tests", "16 assertions against the live gateway"),
    ("35–45", "Policy walkthrough and debug trace", "The six layers, in order, in the portal"),
    ("45–52", "Payer credential and identity options", "SMART, secrets vs certificates, External ID"),
    ("52–58", "Deployment options at Northwind Health", "SKU, network mode, provisioning path"),
    ("58–60", "Decisions, actions, next steps", "Six decisions to own before we leave"),
]
y = 1.50
for i, (when, topic, sub) in enumerate(rows):
    hl = i == 3
    card(s, 0.62, y, 12.1, 0.54, fill=GREEN_LT if hl else WHITE,
         line=GREEN if hl else LINE)
    txt(s, 0.86, y + 0.16, 1.0, 0.28, when, size=12, bold=True,
        color=GREEN if hl else TEAL, font=MONO)
    txt(s, 2.0, y + 0.13, 5.0, 0.3, topic, size=13.5, bold=True, color=INK)
    txt(s, 7.2, y + 0.16, 5.3, 0.28, sub, size=11, color=GREY)
    y += 0.60

txt(s, 0.62, 6.38, 12.1, 0.34,
    "Fifteen of the sixty minutes are a live run against the deployed environment. Everything else supports it.",
    size=13, bold=True, color=TEAL, font=HEAD)
notes(s, """
SAY:
"Quick shape of the hour. Eight topics. The highlighted block — twenty to thirty-five — is the
live demo, and that's the centre of gravity. Everything before it is context so the demo makes
sense, and everything after it is what you'd have to decide to put this into your estate.

Two things I want to flag now. First, at the fifty-eight minute mark there are six decisions
listed. I'd like owners on those before we hang up, not answers necessarily — owners. Second,
if we run long, the segment I'll cut is the policy walkthrough at thirty-five, because that one
is fully written up in the docs and you can read it after. I will not cut the demo."

If asked why no developer portal / monetisation / AI gateway: those are out of scope for today,
happy to book a follow-up.
""")

# ------------------------------------------------------------------ 03 where we left off
s = new()
head(s, "Context", "What was left open on 12 and 18 August")

items = [
    ("1", "Inbound / outbound isolation",
     "\"For inbound and outbound data, can it be isolated from each other within one service? "
     "We don't want payer to query our inbound data.\"", "Platform Engineering, 8/12", GREEN,
     "Answered today — live"),
    ("2", "An APIM demo, with a specialist",
     "\"We can host a separate track on APIM and get an APIM specialist to walk you through "
     "that service and help you get it set up.\"", "Microsoft, 8/18 · 1:15:09", GREEN,
     "This session"),
    ("3", "Offloading authentication to the gateway",
     "\"You may decide to offload the authentication piece … that's where you would do your "
     "private token validation. I know I was using client secrets, but…\"", "Microsoft, 8/18 · 1:15:25",
     AMBER, "Options at 45 min"),
    ("4", "Getting an APIM instance provisioned",
     "\"Let's get that request for APIM going as soon as possible. Even if we have used it, "
     "it is going to take another couple of weeks.\"", "Northwind Health, 8/18 · 1:16:04", AMBER,
     "Action at 58 min"),
]
y = 1.6
for n, title, quote, attrib, col, status in items:
    card(s, 0.62, y, 12.1, 1.24)
    rect(s, 0.62, y, 0.075, 1.24, fill=col)
    chip(s, 0.92, y + 0.2, 0.34, n, fill=col)
    txt(s, 1.44, y + 0.17, 6.4, 0.32, title, size=13.5, bold=True, color=INK)
    txt(s, 1.44, y + 0.53, 9.0, 0.56, quote, size=10.5, color=GREY, italic=True,
        line_spacing=1.2)
    txt(s, 1.44, y + 1.0, 6.0, 0.24, attrib, size=9, bold=True, color=TEAL, font=MONO)
    rect(s, 10.62, y + 0.42, 1.9, 0.4, fill=GREEN_LT if col == GREEN else AMBER_LT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
    txt(s, 10.62, y + 0.52, 1.9, 0.26, status, size=9.5, bold=True, color=col,
        align=PP_ALIGN.CENTER)
    y += 1.34
notes(s, """
SAY:
"Four things were left hanging across the twelfth and the eighteenth. I want to be explicit
about which ones I'm closing today and which ones I'm not.

One and two go green today. The inbound/outbound question gets answered live, and this session
is the specialist walkthrough that was promised.

Three and four are amber deliberately. The credential question — moving off client secrets — I'm
going to give you options at the forty-five minute mark, but the decision is yours and it
depends on what your payers can actually support. And four is a provisioning action, not a
technical one. We'll close that at the end.

I'm quoting you back to yourselves here on purpose. These are from the transcripts, with
timestamps, so if I've misread the intent, stop me now rather than at the end."
""")

# ------------------------------------------------------------------ 04 the decision
s = new()
head(s, "The model", "Physical where a mistake is fatal, logical where it is recoverable")

card(s, 0.62, 1.6, 5.9, 3.5, fill=WHITE, line=GREEN)
rect(s, 0.62, 1.6, 5.9, 0.1, fill=GREEN)
txt(s, 0.94, 1.86, 5.3, 0.34, "PHYSICAL  ·  per payer", size=12.5, bold=True, color=GREEN, font=MONO)
txt(s, 0.94, 2.28, 5.3, 0.5, "One FHIR service per payer", size=19, bold=True, color=INK, font=HEAD)
txt(s, 0.94, 2.86, 5.3, 1.5,
    "Fabrikam's credential has no route to Contoso's FHIR service at all. There is no policy "
    "to get wrong, because there is no path. A defect in APIM cannot cross this line.",
    size=12, color=GREY, line_spacing=1.32)
rect(s, 0.94, 4.36, 5.26, 0.014, fill=LINE)
txt(s, 0.94, 4.54, 5.3, 0.42,
    "Stops a mistake.", size=13, bold=True, color=GREEN, font=HEAD)

card(s, 6.82, 1.6, 5.9, 3.5, fill=WHITE, line=AMBER)
rect(s, 6.82, 1.6, 5.9, 0.1, fill=AMBER)
txt(s, 7.14, 1.86, 5.3, 0.34, "LOGICAL  ·  per contract", size=12.5, bold=True, color=AMBER, font=MONO)
txt(s, 7.14, 2.28, 5.3, 0.5, "meta.tag, enforced by policy", size=19, bold=True, color=INK, font=HEAD)
txt(s, 7.14, 2.86, 5.3, 1.5,
    "Contracts inside one payer are separated by a forced _tag filter. A policy defect here "
    "could expose one contract to another — inside the same payer, who already has a "
    "relationship with you.",
    size=12, color=GREY, line_spacing=1.32)
rect(s, 7.14, 4.36, 5.26, 0.014, fill=LINE)
txt(s, 7.14, 4.54, 5.3, 0.42,
    "Stops a query.", size=13, bold=True, color=AMBER, font=HEAD)

card(s, 0.62, 5.32, 12.1, 1.14, fill=RGBColor(0x0B, 0x22, 0x3A), line=None)
txt(s, 0.96, 5.56, 11.5, 0.66,
    [("Why not a FHIR service per contract?  ", {"bold": True, "color": TEAL_LT}),
     ("Because 30–40 payers × several contracts each is hundreds of instances to operate, "
      "patch and pay for — against a default service limit of 10. Put the hard boundary where "
      "the consequence of failure is highest, and scope the rest at the gateway.",
      {"color": ICE})], size=11.5, line_spacing=1.3)
notes(s, """
This is the conceptual core. If they only remember one slide, make it this one.

SAY:
"This is the decision that was made on the twelfth, and I want to restate it because everything
in the demo follows from it.

Separation comes in two flavours and they are not interchangeable. Physical separation — a
whole FHIR service per payer — is what stops a mistake. If I write a bad policy tomorrow,
Fabrikam still cannot see Contoso's data, because Fabrikam's credential has no route to that
service. There's nothing to get wrong.

Logical separation — contracts inside one payer, tagged and filtered — is what stops a query.
It's enforced by policy, and policy is code, and code can have defects. So the question isn't
'which is better', it's 'where do you put the expensive one'.

The answer we landed on: put the physical boundary at the payer, because that's where a leak
is a breach notification and a headline. Inside a payer, between contracts, a leak is still bad
but it's to a counterparty you already have an agreement with.

And the pragmatic reason — bottom of the slide — the default limit is ten FHIR services per
subscription. You're going to thirty or forty payers. A service per contract is hundreds of
instances. That's not an architecture, that's an operations problem."

If challenged on "is logical good enough for PHI": agree it's a real question, point to test 13b
in the demo — the response body is inspected, not just the status code.
""")

# ------------------------------------------------------------------ 05 architecture
s = new()
head(s, "Architecture", "Where the gateway sits")

def lane(y, label, col):
    rect(s, 0.62, y, 12.1, 0.03, fill=LINE)
    txt(s, 0.62, y + 0.1, 1.5, 0.24, label, size=9, bold=True, color=col, font=MONO)

txt(s, 0.62, 1.5, 2.4, 0.28, "PAYERS", size=10, bold=True, color=TEAL, font=MONO)
for i, (nm, ct) in enumerate([("Contoso Health", "CT-3456 · CT-7788"), ("Fabrikam Plan", "CT-9012")]):
    card(s, 0.62, 1.82 + i * 1.0, 2.5, 0.84)
    txt(s, 0.86, 1.98 + i * 1.0, 2.1, 0.3, nm, size=12, bold=True, color=INK)
    txt(s, 0.86, 2.3 + i * 1.0, 2.1, 0.24, ct, size=9.5, color=GREY, font=MONO)

txt(s, 0.62, 3.96, 2.5, 0.9,
    "SMART Backend Services\nclient_credentials\nEntra-issued token",
    size=10, color=GREY, font=MONO, line_spacing=1.35)

txt(s, 3.5, 1.5, 2.4, 0.28, "GATEWAY", size=10, bold=True, color=TEAL, font=MONO)
card(s, 3.5, 1.82, 3.2, 3.5, fill=RGBColor(0x0B, 0x22, 0x3A), line=None)
txt(s, 3.74, 2.02, 2.8, 0.3, APIM, size=9.5, bold=True, color=TEAL_LT, font=MONO)
layers = ["1  validate-jwt", "2  entitlement", "3  route allow-list",
          "4  scope  _tag / Group", "5  rate limit", "6  managed identity"]
for i, l in enumerate(layers):
    rect(s, 3.74, 2.44 + i * 0.44, 2.72, 0.36, fill=RGBColor(0x14, 0x33, 0x54),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    txt(s, 3.9, 2.52 + i * 0.44, 2.5, 0.24, l, size=10, color=ICE, font=MONO)

txt(s, 7.06, 1.5, 2.4, 0.28, "ROUTES", size=10, bold=True, color=TEAL, font=MONO)
card(s, 7.06, 1.82, 2.6, 1.5, fill=GREEN_LT, line=GREEN)
txt(s, 7.28, 2.0, 2.2, 0.28, "OUTBOUND", size=11, bold=True, color=GREEN, font=MONO)
txt(s, 7.28, 2.32, 2.2, 0.9, "GET · search\nGroup/{id}/$export\nno writes", size=10,
    color=INK, font=MONO, line_spacing=1.3)
card(s, 7.06, 3.5, 2.6, 1.5, fill=AMBER_LT, line=AMBER)
txt(s, 7.28, 3.68, 2.2, 0.28, "INBOUND", size=11, bold=True, color=AMBER, font=MONO)
txt(s, 7.28, 4.0, 2.2, 0.9, "POST · PUT\n$import\nno $export", size=10,
    color=INK, font=MONO, line_spacing=1.3)

txt(s, 10.06, 1.5, 2.6, 0.28, "DATA PLANE", size=10, bold=True, color=TEAL, font=MONO)
for i, nm in enumerate(["ahdspocdemo01\n-fhir-payera", "ahdspocdemo01\n-fhir-payerb"]):
    card(s, 10.06, 1.82 + i * 1.62, 2.66, 1.34)
    rect(s, 10.06, 1.82 + i * 1.62, 2.66, 0.08, fill=TEAL)
    txt(s, 10.3, 2.04 + i * 1.62, 2.2, 0.5, nm, size=9.5, bold=True, color=INK,
        font=MONO, line_spacing=1.25)
    txt(s, 10.3, 2.62 + i * 1.62, 2.2, 0.4, "FHIR Data Contributor\ngranted to APIM only",
        size=8.5, color=GREY, font=MONO, line_spacing=1.25)

card(s, 0.62, 5.56, 12.1, 0.94, fill=RED_LT, line=RED)
txt(s, 0.96, 5.76, 11.4, 0.56,
    [("The payer's Entra app holds no Azure RBAC role on either FHIR service.  ", {"bold": True, "color": RED}),
     ("A payer token sent directly to AHDS returns 403 — not 401. Entra issued it; it simply "
      "carries no authority. That is what makes the gateway a boundary rather than a speed bump.",
      {"color": INK})], size=11.5, line_spacing=1.3)
notes(s, """
Walk this left to right. Do not rush — this is the mental model the demo assumes.

SAY:
"Left to right. Two payers. Contoso holds two contracts, Fabrikam holds one. They authenticate
with SMART Backend Services — client credentials against Entra — so no human, no browser, no
refresh token. This is server-to-server.

Middle: the gateway. Six policy layers, and I'll come back to the ordering because the order is
load-bearing, not cosmetic.

Then routes. This is the answer to the eighteenth. Each payer gets two APIs, not one. Outbound
is read-only — search, and Group-scoped export. Inbound accepts writes and import and explicitly
refuses export. Same FHIR service behind both.

Right side: the data plane. One FHIR service per payer.

Now the red box, and this is the part I'd underline. The payer's application has no Azure role
on the FHIR service. None. Only APIM's managed identity holds FHIR Data Contributor. So if a
payer takes their perfectly valid token and points it straight at the AHDS URL, bypassing us,
they get a 403. Not a 401 — the token is real. It just carries no authority.

That's the difference between filtering and a boundary. I'll prove that one live, it's test
twelve."
""")

# ------------------------------------------------------------------ 06 the answer
s = new()
head(s, "The question from 18 August", "One FHIR service, two routes, two credentials")

txt(s, 0.62, 1.5, 12.1, 0.5,
    "\"For inbound and outbound data, can it be isolated from each other within one service? "
    "We don't want payer to query our inbound data.\"",
    size=14.5, italic=True, color=INK, font=HEAD, line_spacing=1.25)
txt(s, 0.62, 2.08, 12.1, 0.26, "PLATFORM ENGINEERING · 12 AUGUST", size=9.5, bold=True,
    color=TEAL, font=MONO)

card(s, 0.62, 2.62, 12.1, 0.72, fill=RGBColor(0x0B, 0x22, 0x3A), line=None)
txt(s, 0.96, 2.82, 11.5, 0.36,
    [("Yes. Isolation is not a property of the FHIR service — it is a property of the route.",
      {"bold": True, "color": ICE})], size=14, font=HEAD)

hdrs = ["", "Outbound API  ·  payera/outbound", "Inbound API  ·  payera/inbound"]
rows2 = [
    ("Credential", "Payer's Entra app", "Northwind Health ingest principal"),
    ("GET / search", "allowed, _tag forced", "allowed"),
    ("POST / PUT", "403 denied", "allowed, contract header required"),
    ("Group/{id}/$export", "allowed, own Groups only", "403 denied"),
    ("System / Patient $export", "403 denied", "403 denied"),
    ("Rate limit", "600/60s  ·  export 1 per 300s", "3000/60s"),
]
x0, colw = 0.62, [3.0, 4.55, 4.55]
y = 3.4
rect(s, x0, y, 12.1, 0.42, fill=NAVY_CARD)
cx = x0
for i, h in enumerate(hdrs):
    txt(s, cx + 0.18, y + 0.1, colw[i] - 0.3, 0.28, h, size=10.5, bold=True,
        color=WHITE, font=MONO)
    cx += colw[i]
y += 0.42
for r_i, row in enumerate(rows2):
    fill = WHITE if r_i % 2 == 0 else RGBColor(0xEE, 0xF3, 0xF7)
    rect(s, x0, y, 12.1, 0.42, fill=fill)
    cx = x0
    for i, cell in enumerate(row):
        col = INK
        bold = i == 0
        if "403" in cell:
            col = RED
            bold = True
        elif "allowed" in cell:
            col = GREEN
        txt(s, cx + 0.18, y + 0.1, colw[i] - 0.3, 0.28, cell, size=10.5, color=col,
            bold=bold, font=BODY if i == 0 else MONO)
        cx += colw[i]
    y += 0.42

txt(s, 0.62, 6.46, 12.1, 0.34,
    "Deny by default: any FHIR operation not on the allow-list is refused, not proxied. New operations age safely.",
    size=12, bold=True, color=TEAL, font=HEAD)
notes(s, """
This is THE slide for Ashay's question. Slow down. Ask for confirmation before moving on.

SAY:
"Here's the question as it was asked on the twelfth, and here's the answer.

Yes — and the reason it works is worth internalising. Isolation is not a property of the FHIR
service. It's a property of the route. Same database behind both columns. What differs is which
API you're calling, which credential opens it, and what that API's policy permits.

Read the table down. Different credential — the payer's app on the left, your ingest principal
on the right, and neither one is onboarded to the other's API. Writes are refused on outbound.
Export is refused on inbound, which matters because it means a compromised ingest credential
cannot be turned into a bulk exfiltration tool. And system-level and patient-level export are
refused on both, because they're unbounded — the working set is the entire instance.

Bottom line: deny by default. If a FHIR operation isn't on the allow-list, it's refused rather
than passed through. That's deliberate — FHIR gains operations over time, and an allow-list ages
safely where a deny-list does not."

THEN ASK: "Ashay — is that the isolation you were asking about, or were you picturing something
stronger, like separate services for inbound and outbound?" Let them answer before moving on.
""")

# ------------------------------------------------------------------ 07 six layers
s = new(dark=True)
head(s, "Policy design", "Six layers, ordered by cost", dark=True)

layers = [
    ("1", "authenticate", "Is this token real?", "crypto, cached JWKS", TEAL_LT),
    ("2", "entitle", "Is this caller onboarded to this payer?", "named-value lookup", TEAL_LT),
    ("3", "allow-list route", "Is this verb + path permitted?", "string match", TEAL_LT),
    ("4", "scope", "Narrow the request to what they hold", "query rewrite", TEAL_LT),
    ("5", "rate limit", "Is this caller within budget?", "counter", AMBER),
    ("6", "broker", "Swap in the managed identity", "token acquisition — network call", GREEN),
]
y = 1.66
for n, name, q, cost, col in layers:
    card(s, 0.62, y, 12.1, 0.68, fill=NAVY_CARD, line=SLATE)
    chip(s, 0.9, y + 0.16, 0.36, n, fill=col, color=NAVY if col != TEAL_LT else NAVY)
    txt(s, 1.5, y + 0.2, 2.5, 0.3, name, size=13, bold=True, color=WHITE, font=MONO)
    txt(s, 4.2, y + 0.21, 5.4, 0.3, q, size=12, color=ICE)
    txt(s, 9.9, y + 0.23, 2.6, 0.26, cost, size=10, color=MUTED, font=MONO,
        align=PP_ALIGN.RIGHT)
    y += 0.76

card(s, 0.62, 6.28, 12.1, 0.5, fill=RGBColor(0x0B, 0x22, 0x3A), line=SLATE)
txt(s, 0.94, 6.4, 11.5, 0.28,
    [("Each layer is cheap relative to the next. Reject as early as possible — and never spend a "
      "token acquisition on a request that was going to be refused.", {"color": TEAL_LT})],
    size=11.5, bold=True)
notes(s, """
SAY:
"Six layers. The content matters less than the order, so let me explain the ordering.

Each layer is cheaper than the one below it. Validating a signature is cryptography against a
cached key — microseconds. A named-value lookup is memory. A route check is a string match.
Rejecting early costs you almost nothing.

Layer six is last on purpose. Acquiring a managed identity token is a network call. If you did
that first, every unauthorised request would burn a token acquisition before you refused it.
That's how you turn a security control into a denial-of-service amplifier.

Layer five — rate limiting — sits after layer three, and this one is subtler. If you rate-limit
before you authorise, then rejected requests eat the caller's quota. A payer sending malformed
requests would exhaust their own budget and then legitimate traffic fails. Count only what you
were actually going to serve.

I'm making a point of this because when you take these policies and adapt them, the temptation
is to reorder for readability. Don't. The order is a design decision."
""")

# ------------------------------------------------------------------ 08 auth + entitle
s = new()
head(s, "Layers 1–2", "Authenticate, then entitle")

txt(s, 0.62, 1.5, 5.9, 0.34, "1  ·  VALIDATE-JWT", size=12, bold=True, color=TEAL, font=MONO)
code(s, 0.62, 1.86, 5.9, 2.0, [
    "<validate-jwt header-name=\"Authorization\"",
    "   failed-validation-httpcode=\"401\"",
    "   output-token-variable-name=\"jwt\">",
    "  <openid-config url=\"…/v2.0/.well-known/…\" />",
    ("  <audiences>{fhirServiceUrl}</audiences>", GREEN_LT),
    "  <issuers>sts.windows.net/… , login…/v2.0</issuers>",
    "</validate-jwt>",
], size=9.5)
txt(s, 0.62, 4.0, 5.9, 1.7,
    [("The audience is the FHIR service URL, not the gateway.", {"bold": True, "color": INK}),
     ("A token minted for Payer B fails on Payer A's route before any entitlement logic runs — "
      "two independent gates, cheapest first.", {"color": GREY, "space_before": 5}),
     ("Both issuer forms are listed because Entra emits v1 or v2 depending on the app "
      "registration. Listing one causes intermittent 401s that look payer-specific.",
      {"color": GREY, "space_before": 7})],
    size=11, line_spacing=1.3)

txt(s, 6.82, 1.5, 5.9, 0.34, "2  ·  ENTITLEMENT", size=12, bold=True, color=TEAL, font=MONO)
code(s, 6.82, 1.86, 5.9, 2.0, [
    "// APIM named value: payer-entitlements",
    "{ \"<appId>\": {",
    ("    \"payer\":     \"payera\",", GREEN_LT),
    "    \"contracts\": [\"CT-3456\",\"CT-7788\"],",
    "    \"groups\":    [\"group-ct3456\", …] } }",
    "",
    "// azp (v2) or appid (v1) → the caller",
], size=9.5)
txt(s, 6.82, 4.0, 5.9, 1.7,
    [("The cross-payer guard.", {"bold": True, "color": INK}),
     ("If entitlement.payer does not match the route's payer, refuse by name — 403. A valid "
      "Payer B credential is rejected on Payer A's route even if it somehow validated.",
      {"color": GREY, "space_before": 5}),
     ("At 30–40 payers this named value becomes a cached lookup against the contract master. "
      "The policy shape does not change — only the source does.", {"color": GREY, "space_before": 7})],
    size=11, line_spacing=1.3)

card(s, 0.62, 5.74, 12.1, 0.86, fill=AMBER_LT, line=AMBER)
txt(s, 0.96, 5.92, 11.4, 0.3,
    [("Editing gotcha:  ", {"bold": True, "color": AMBER}),
     ("the entitlement JSON is stored single-quoted and swapped to double quotes at parse time. "
      "JSON inside a C# literal inside XML needs escaping at two levels — that is where these "
      "policies usually break.", {"color": INK})], size=11)
notes(s, """
Only go into this depth if the room is technical and engaged. Otherwise summarise in 30 seconds
and move on — it's all written up in docs/07-apim-control-plane.md.

SAY:
"Two layers, side by side.

On the left, token validation. The one detail I'd point at is the audience. We set it to the
FHIR service URL, not the gateway URL. That means a token minted for Payer B's service fails
validation on Payer A's route immediately — before we do any entitlement work at all. Two
independent gates and the cheap one runs first.

The issuer line is a scar. Entra emits either a v1 or a v2 issuer depending on how the app
registration is configured. If you list only one, you get intermittent 401s that appear to
depend on which payer is calling — because they do. List both.

On the right, entitlement. This resolves the calling application to the contracts it may see.
The important line is 'payer'. If the entitlement says payera and the route is payerb, we refuse
by name. That's the cross-payer guard, and it's what makes physical separation meaningful rather
than decorative.

For the POC this is a named value — visible in the portal, auditable, no dependency. At thirty
or forty payers you'd make it a cached lookup against your contract master. The policy shape
doesn't change, only where the record comes from.

Bottom box — practical warning. That JSON lives inside a C# expression inside XML. Quotes need
escaping at two levels. We store it single-quoted and swap at parse time. This is the single
most common way these policies break during editing."
""")

# ------------------------------------------------------------------ 09 scoping
s = new()
head(s, "Layer 4", "Scoping — the contract boundary")

txt(s, 0.62, 1.5, 5.9, 0.34, "GROUP ENTITLEMENT", size=12, bold=True, color=TEAL, font=MONO)
card(s, 0.62, 1.86, 5.9, 1.5)
txt(s, 0.94, 2.1, 5.3, 0.34, "Group/group-ct3456/$export", size=14, bold=True,
    color=INK, font=MONO)
rect(s, 2.36, 2.5, 1.72, 0.03, fill=AMBER)
txt(s, 0.94, 2.62, 5.3, 0.6,
    "must appear in entitlement.groups[]", size=11, color=AMBER, font=MONO, bold=True)
txt(s, 0.94, 2.94, 5.3, 0.36,
    "Right payer, valid token — still cannot export a cohort it does not hold.",
    size=11, color=GREY)

txt(s, 6.82, 1.5, 5.9, 0.34, "FORCED  _tag", size=12, bold=True, color=TEAL, font=MONO)
code(s, 6.82, 1.86, 5.9, 1.5, [
    "<set-query-parameter name=\"_tag\"",
    ("   exists-action=\"override\">", GREEN_LT),
    "  <value>…/contract|CT-3456,…|CT-7788</value>",
    "</set-query-parameter>",
], size=10)

card(s, 0.62, 3.62, 12.1, 1.16, fill=GREEN_LT, line=GREEN)
txt(s, 0.96, 3.84, 11.4, 0.72,
    [("exists-action=\"override\" is the load-bearing attribute.", {"bold": True, "color": GREEN}),
     ("A caller-supplied _tag is replaced, not merged. A payer cannot widen its own result set "
      "by asking for a contract it does not hold. This is assertion 13 in the demo — and 13b "
      "inspects the response body, not just the status code.", {"color": INK, "space_before": 5})],
    size=12, line_spacing=1.3)

txt(s, 0.62, 5.06, 5.9, 0.34, "AND ON THE WAY IN", size=12, bold=True, color=TEAL, font=MONO)
txt(s, 0.62, 5.42, 12.1, 1.1,
    [("The inbound route stamps every write with its contract tag and refuses writes that "
      "arrive without an X-Payer-Contract header.", {"color": INK}),
     ("If ingest does not stamp, export cannot filter. Refusing untagged writes is deliberate: "
      "silently untagged data is invisible to every payer, which means it is effectively lost — "
      "and you would not find out until someone asked for it.",
      {"color": GREY, "space_before": 6})],
    size=12, line_spacing=1.32)
notes(s, """
SAY:
"Layer four is where contract-level separation actually happens, and there are two halves to it.

Top left — Group entitlement. When a payer asks to export a cohort, we check the Group ID
against the list in their entitlement record. So even inside the right payer, with a valid
token, on the right route, they cannot export a cohort that isn't theirs.

Top right — the forced tag. Every search gets a _tag query parameter injected, scoped to the
contracts the caller holds. And the attribute that matters is exists-action equals override.
Override, not merge. If the payer supplies their own _tag asking for a contract they don't hold,
we throw theirs away and substitute ours. They cannot widen their own result set. That's
assertion thirteen in the demo, and thirteen-b actually opens the response body and checks that
only their contracts came back — not just that the status code was 200.

Bottom half, and this is the piece people miss. None of that outbound filtering works unless
ingest stamps the data on the way in. So the inbound route requires a contract header on every
write and refuses the write without it.

That refusal is deliberate and I want to defend it, because someone will ask why we don't just
default it. If you silently accept untagged data, that data is invisible to every payer's tag
filter. It's in the database and nobody can ever see it. You'd discover that months later when
a payer asks why a member's claims are missing. Fail loudly at write time instead."
""")

# ------------------------------------------------------------------ 10 rate limiting
s = new()
head(s, "Layer 5", "Rate limiting — the burst mitigation")

txt(s, 0.62, 1.5, 12.1, 0.44,
    "\"They usually happen at the same time … it will pop the server itself.\"",
    size=15, italic=True, color=INK, font=HEAD)
txt(s, 0.62, 2.0, 12.1, 0.26, "ON PAYERS SUBMITTING BULK EXPORTS SIMULTANEOUSLY",
    size=9.5, bold=True, color=TEAL, font=MONO)

lims = [
    ("600 / 60s", "per payer", "Normal request protection", TEAL),
    ("50,000 / day", "per payer", "Abuse ceiling", TEAL),
    ("1 / 300s", "per payer", "$export submissions only", AMBER),
]
for i, (v, unit, desc, col) in enumerate(lims):
    x = 0.62 + i * 4.06
    card(s, x, 2.52, 3.86, 1.5, fill=AMBER_LT if col == AMBER else WHITE,
         line=col if col == AMBER else LINE)
    txt(s, x + 0.3, 2.76, 3.3, 0.5, v, size=25, bold=True, color=col, font=HEAD)
    txt(s, x + 0.3, 3.26, 3.3, 0.26, unit, size=10.5, color=GREY, font=MONO)
    txt(s, x + 0.3, 3.56, 3.3, 0.36, desc, size=11, color=INK)

code(s, 0.62, 4.22, 12.1, 1.2, [
    "<rate-limit-by-key calls=\"1\" renewal-period=\"300\"",
    "   counter-key=\"@(\"export-\" + payerKey)\"",
    ("   increment-condition=\"@(routeClass == \"export\")\" />", GREEN_LT),
], size=11.5)

card(s, 0.62, 5.6, 12.1, 0.94, fill=GREEN_LT, line=GREEN)
txt(s, 0.96, 5.8, 11.4, 0.56,
    [("Only submissions increment the counter.  ", {"bold": True, "color": GREEN}),
     ("Polling an in-flight job is unaffected, so a payer is never locked out of the job they "
      "already started. The burst becomes a queue at the gateway — before any FHIR capacity is "
      "consumed.", {"color": INK})], size=11.5, line_spacing=1.3)
notes(s, """
For Northwind Health this is arguably the highest-value single line in the whole policy set. Sell it.

SAY:
"This one came straight out of your own concern on the earlier call — that payers tend to submit
their bulk exports at the same time, month end or quarter end, and that the coincidence is what
takes the server down.

Three limits. The first two are ordinary hygiene. The third is the interesting one: one export
submission per payer per five minutes.

Look at the increment-condition on the policy. The counter only ticks when a payer *submits* an
export. Polling a job that's already running doesn't touch it. That matters — otherwise a payer
would start a job, poll it, and lock themselves out of their own job. That's the kind of bug
that generates a support call at two in the morning.

What this does architecturally: it converts a simultaneous burst into a queue, and the queue
forms at the gateway, before any FHIR capacity is consumed. AHDS never sees the stampede.

It's assertion fifteen in the demo — second export inside five minutes returns 429. And it's the
reason I have to warm the environment before this call, because the lock is held per payer for
three hundred seconds."
""")

# ------------------------------------------------------------------ 11 broker
s = new(dark=True)
head(s, "Layer 6", "Trusted broker — why the gateway cannot be bypassed", dark=True)

code(s, 0.62, 1.6, 12.1, 1.16, [
    "<authentication-managed-identity resource=\"{fhirUrl}\" output-token-variable-name=\"msi\" ignore-error=\"false\" />",
    "<set-header name=\"Authorization\" exists-action=\"override\">",
    ("  <value>@(\"Bearer \" + (string)context.Variables[\"msi\"])</value>", GREEN_LT),
], size=11, fill=RGBColor(0x08, 0x1B, 0x30))

txt(s, 0.62, 2.96, 12.1, 0.36,
    "The payer's token is discarded. AHDS is called with APIM's managed identity.",
    size=15, bold=True, color=ICE, font=HEAD)

cons = [
    ("The payer credential never reaches the FHIR service", GREEN),
    ("AHDS makes no payer-specific authorisation decision — it sees one trusted caller", GREEN),
    ("Only APIM holds FHIR Data Contributor, so a direct payer call returns 403", GREEN),
    ("ignore-error=\"false\" — if the MI token fails, refuse rather than forward unauthenticated", GREEN),
]
y = 3.5
for t, col in cons:
    card(s, 0.62, y, 12.1, 0.54, fill=NAVY_CARD, line=SLATE)
    chip(s, 0.9, y + 0.13, 0.28, "✓", fill=col, color=WHITE, size=11)
    txt(s, 1.42, y + 0.15, 11.0, 0.3, t, size=12, color=ICE)
    y += 0.62

card(s, 0.62, 6.06, 12.1, 0.64, fill=RGBColor(0x08, 0x1B, 0x30), line=TEAL)
txt(s, 0.96, 6.24, 11.4, 0.3,
    [("Audit is preserved:  ", {"bold": True, "color": TEAL_LT}),
     ("X-MS-AZUREFHIR-AUDIT-CALLER, -PAYER and -CONTRACTS are added so the AHDS audit log can "
      "still attribute a request to the originating payer despite the identity swap.", {"color": ICE})],
    size=11.5)
notes(s, """
This is the keystone slide for the security audience. Land it firmly.

SAY:
"Last layer, and this is the keystone.

Once every check has passed, we throw the payer's token away. We call AHDS with API Management's
own managed identity.

Four consequences, all of them good. The payer's credential never touches the FHIR service.
AHDS isn't making payer-specific authorisation decisions — from its point of view there's one
trusted caller. Only APIM holds the FHIR data role. And if the managed identity token can't be
acquired, we fail the request rather than forwarding it unauthenticated — that's the
ignore-error equals false.

The third bullet is the one that matters for your security review. It's what turns 'we filter at
the gateway' into 'the gateway cannot be bypassed'. If a payer discovers the real AHDS hostname
— and they might, hostnames leak — and calls it directly with their valid token, they get a 403.
Because their app has no role. There is no configuration they can reach that changes that.

One thing people worry about: does the identity swap destroy your audit trail? No. We add two
headers carrying the original caller and payer, so AHDS audit logs still attribute the request
correctly."
""")

# ------------------------------------------------------------------ 12 demo
s = new(dark=True)
head(s, "Live demo", "Sixteen assertions against the deployed gateway", dark=True)

txt(s, 0.62, 1.52, 7.4, 0.4, "./scripts/run-isolation-tests.ps1", size=17, bold=True,
    color=TEAL_LT, font=MONO)
txt(s, 0.62, 1.98, 7.4, 0.9,
    "Mints a short-lived credential for each of two payers, exercises sixteen assertions "
    "against the live gateway, then revokes them. Nothing is written to disk.",
    size=11.5, color=MUTED, line_spacing=1.3)

tests = [
    ("1–3", "own data readable · Group export accepted · capability statement", GREEN),
    ("4", "payer B app with valid payer A audience  →  403", RED),
    ("5", "payer B token, payer B audience, payer A route  →  401", RED),
    ("6", "unentitled Group export  →  403", RED),
    ("7", "write on the outbound route  →  403", RED),
    ("8–9", "payer credential on inbound route · export on inbound  →  403", RED),
    ("10–11", "system-level and patient-level export  →  403", RED),
    ("12", "payer token sent straight to AHDS  →  403", RED),
    ("13", "caller-supplied _tag overridden  ·  13b body carries only own contracts", GREEN),
    ("14", "untagged inbound write  →  rejected", RED),
    ("15", "second export within five minutes  →  429", AMBER),
]
y = 3.02
for n, t, col in tests:
    rect(s, 0.62, y, 8.9, 0.3, fill=None)
    txt(s, 0.62, y, 0.8, 0.26, n, size=10.5, bold=True, color=col, font=MONO)
    txt(s, 1.52, y, 8.0, 0.26, t, size=10.5, color=ICE, font=MONO)
    y += 0.325

card(s, 9.84, 1.52, 2.88, 4.9, fill=RGBColor(0x08, 0x1B, 0x30), line=TEAL)
txt(s, 10.1, 1.82, 2.4, 0.3, "POINT AT THESE", size=10, bold=True, color=TEAL_LT, font=MONO)
picks = [
    ("4", "A legitimate token from the wrong payer is refused. The PHI boundary holds even when the credential is real."),
    ("12", "403, not 401. Entra issued the token; the app holds no role. Boundary, not speed bump."),
    ("13b", "The response body was inspected, not just the status code."),
]
yy = 2.24
for n, why in picks:
    chip(s, 10.1, yy, 0.34, n, fill=TEAL, size=11)
    txt(s, 10.1, yy + 0.44, 2.4, 1.1, why, size=10, color=ICE, line_spacing=1.28)
    yy += 1.44
notes(s, """
=== THIS IS THE CENTRE OF THE SESSION — 15 MINUTES ===

BEFORE THE CALL: run the suite once to warm the gateway, then wait 5 minutes (the export lock is
300s per payer) or pass -SkipThrottleTest.

Switch to the terminal. Run it. While it runs (~3 min), say:

"What this is doing: it mints a short-lived credential for each of two payers, runs sixteen
assertions against the live gateway, and revokes the credentials at the end. Nothing is written
to disk, and I'm not using any pre-recorded output — if the model is wrong this goes red while
you're watching."

WHEN THE TABLE APPEARS, do not read all sixteen. Point at three:

"Line four. Payer B's application, presenting a token with the correct audience for Payer A's
service. Everything about that token is legitimate. Refused. That's the PHI boundary, and it
holds even when the credential itself is valid.

Line twelve is the one I'd take to your security review. That's the payer's token sent directly
to the AHDS endpoint, bypassing the gateway entirely. It returns 403, not 401. The token is
real — Entra minted it. The application simply holds no role on that FHIR service. That's why
this is a boundary and not a speed bump.

Line thirteen-b. We didn't just check the status code, we opened the response body and confirmed
only the caller's own contracts came back. Status codes tell you a request was refused; bodies
tell you a filter actually worked."

THEN OFFER: "Anything you want me to break? I can add a case and re-run — that's a more useful
use of the next five minutes than me talking."

If a test fails live: do not hide it. Say "that's a real failure, let me capture it and come
back to you" and move on. Credibility is the asset here.
""")

# ------------------------------------------------------------------ 13 credentials
s = new()
head(s, "Credential model", "What payers present, and where they live")

txt(s, 0.62, 1.5, 12.1, 0.44,
    "\"You may decide to offload the authentication piece … that's where you would do your "
    "private token validation. I know I was using client secrets, but…\"",
    size=13, italic=True, color=INK, line_spacing=1.25)
txt(s, 0.62, 2.02, 12.1, 0.26, "MICROSOFT · 18 AUGUST · 1:15:25", size=9.5, bold=True,
    color=TEAL, font=MONO)

opts = [
    ("Client secret", "POC only", RED_LT, RED,
     "What the POC uses. Shared secret, rotation is manual, and it is copyable — anyone who "
     "reads it once can use it forever."),
    ("private_key_jwt", "SMART standard", GREEN_LT, GREEN,
     "The payer holds a private key and signs its own assertion. Nothing shared and nothing "
     "copyable. This is what SMART Backend Services actually specifies."),
    ("mTLS client cert", "Gateway-enforced", GREEN_LT, GREEN,
     "APIM validates a client certificate at the edge, before any policy runs. Can be combined "
     "with either of the above as a second factor."),
]
for i, (name, tag, bg, col, desc) in enumerate(opts):
    x = 0.62 + i * 4.06
    card(s, x, 2.5, 3.86, 2.1, fill=WHITE, line=col)
    rect(s, x, 2.5, 3.86, 0.09, fill=col)
    txt(s, x + 0.28, 2.74, 3.3, 0.34, name, size=14, bold=True, color=INK, font=MONO)
    rect(s, x + 0.28, 3.14, 1.5, 0.3, fill=bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    txt(s, x + 0.28, 3.2, 1.5, 0.22, tag, size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, x + 0.28, 3.6, 3.3, 0.9, desc, size=10.5, color=GREY, line_spacing=1.28)

txt(s, 0.62, 4.86, 12.1, 0.34, "WHERE THE PAYER IDENTITY LIVES", size=12, bold=True,
    color=TEAL, font=MONO)
card(s, 0.62, 5.24, 5.9, 1.24, fill=GREEN_LT, line=GREEN)
txt(s, 0.94, 5.44, 5.3, 0.3, "External ID tenant", size=13, bold=True, color=GREEN)
txt(s, 0.94, 5.78, 5.3, 0.6,
    "A separate tenant hanging off yours. Payers never become objects in the Northwind Health "
    "enterprise directory.", size=11, color=INK, line_spacing=1.28)
card(s, 6.82, 5.24, 5.9, 1.24, fill=WHITE, line=LINE)
txt(s, 7.14, 5.44, 5.3, 0.3, "Northwind Health enterprise tenant", size=13, bold=True, color=GREY)
txt(s, 7.14, 5.78, 5.3, 0.6,
    "Simpler — assign RBAC directly. But it means external parties as guests in your corporate "
    "directory. Most organisations decline.", size=11, color=GREY, line_spacing=1.28)
notes(s, """
SAY:
"Steve raised this on the eighteenth and it's a genuine decision, so here are the options rather
than a recommendation.

Three credential types. The POC uses client secrets, and I want to be blunt that this is a POC
convenience, not a recommendation. A shared secret is copyable — anyone who sees it once can use
it indefinitely — and rotation is a manual coordination exercise with every payer.

Private key JWT is what SMART Backend Services actually specifies. The payer holds a private key
and signs its own client assertion. Nothing is shared, so nothing is copyable. If your payers
are mature FHIR implementers they will already support this, because it's the standard.

Mutual TLS is the third option, and it's what Steve was describing — install the certificate on
the APIM instance and validate the client cert at the edge, before any policy code runs. You can
combine it with either of the others as a second factor.

My honest read: go to private_key_jwt because it's the standard and your payers will expect it,
and hold mTLS in reserve for a payer who specifically asks for network-level assurance.

Bottom half — where payer identities live. Strong recommendation for External ID. The
alternative works and is simpler, but it means external organisations become guest objects in
the Northwind Health corporate directory. In my experience almost nobody accepts that once security
reviews it."

ASK: "Do you know yet what your payers can support? That constrains this more than our
preference does."
""")

# ------------------------------------------------------------------ 14 capacity
s = new()
head(s, "Capacity", "Service limits, and what the gateway does about them")

card(s, 0.62, 1.6, 5.9, 2.2, fill=WHITE, line=LINE)
txt(s, 0.94, 1.84, 5.3, 0.3, "FHIR SERVICES PER SUBSCRIPTION", size=10.5, bold=True,
    color=TEAL, font=MONO)
txt(s, 0.94, 2.26, 2.2, 0.7, "10", size=42, bold=True, color=GREY, font=HEAD)
txt(s, 2.6, 2.46, 0.8, 0.5, "→", size=28, bold=True, color=TEAL, font=HEAD)
txt(s, 3.5, 2.26, 2.6, 0.7, "40", size=42, bold=True, color=GREEN, font=HEAD)
txt(s, 0.94, 3.06, 5.3, 0.56,
    "Default, raised by support request. Approved for this POC. One service per payer means "
    "this is the number that governs how many payers you can onboard.",
    size=11, color=GREY, line_spacing=1.28)

card(s, 6.82, 1.6, 5.9, 2.2, fill=WHITE, line=LINE)
txt(s, 7.14, 1.84, 5.3, 0.3, "THE EXPORT PROBLEM", size=10.5, bold=True, color=TEAL, font=MONO)
txt(s, 7.14, 2.24, 5.3, 1.4,
    [("System-level and Patient-level $export are unbounded — the working set is the whole "
      "instance.", {"color": INK}),
     ("Group-scoped export bounds it to a cohort. That is simultaneously a security property "
      "and the capacity mitigation.", {"color": GREY, "space_before": 6})],
    size=11.5, line_spacing=1.3)

card(s, 0.62, 4.0, 12.1, 1.2, fill=AMBER_LT, line=AMBER)
txt(s, 0.96, 4.2, 11.4, 0.8,
    [("Two AHDS behaviours the gateway has to absorb", {"bold": True, "color": AMBER}),
     ("$export returns 400 unless Accept: application/fhir+json is set explicitly — so the "
      "gateway injects it, rather than every payer client having to know.  ·  Content-Location "
      "on a 202 points at the AHDS host; the policy rewrites it to the gateway host, or the "
      "payer learns the real endpoint and polls where the policies do not apply.",
      {"color": INK, "space_before": 5})], size=11.5, line_spacing=1.3)

card(s, 0.62, 5.4, 12.1, 1.06, fill=RGBColor(0x0B, 0x22, 0x3A), line=None)
txt(s, 0.96, 5.6, 11.5, 0.66,
    [("$import 403 — root cause, for the record.  ", {"bold": True, "color": TEAL_LT}),
     ("$import reads the integration storage account as the FHIR service's own system-assigned "
      "identity, not a user-assigned identity attached to the service. The dev environment had "
      "the storage role on the UAMI while the read was being attempted by the service identity, "
      "which held nothing. Fixed and annotated in infra/modules/rbac.bicep.", {"color": ICE})],
    size=11, line_spacing=1.3)
notes(s, """
SAY:
"Three capacity points, quickly.

Top left. The default limit is ten FHIR services per subscription, raised to forty by support
request — already approved for this POC. Since the model is one service per payer, that number
is literally your payer ceiling. Worth knowing before you plan the onboarding schedule.

Top right. Why we only allow Group-scoped export. System-level and patient-level export are
legitimate FHIR, but they're unbounded — the working set is the entire instance. Group-scoped
bounds it to a cohort. That's a nice property because the security control and the capacity
control are the same control.

Amber box — two AHDS quirks the gateway absorbs so your payers never see them. First, export
returns a 400 unless the Accept header is set to application/fhir+json explicitly. Rather than
document that for thirty payers, the gateway injects it. Second, the poll URL that comes back on
a 202 points at the real AHDS hostname. We rewrite it to the gateway host — otherwise the payer
discovers your actual endpoint and starts polling somewhere our policies don't apply.

Bottom, dark box — the import 403 from the earlier call. Root cause: import reads the storage
account as the FHIR service's own system-assigned identity, not as a user-assigned identity you
attach to it. The role was on the wrong principal. It's fixed, and I've annotated the Bicep so
nobody undoes it in a refactor six months from now."
""")

# ------------------------------------------------------------------ 15 landing it
s = new()
head(s, "Deployment", "What it takes to land this at Northwind Health")

cols = [
    ("SKU", "BasicV2 in the POC", [
        "Developer / BasicV2 — no SLA, POC only",
        "StandardV2 — SLA, no VNet injection",
        "Premium — VNet injection, multi-region, zones",
        "Choice is driven by network mode, not throughput",
    ]),
    ("Network", "Decide with the proxy team", [
        "External — gateway on a public IP",
        "Internal VNet — gateway private, fronted by your edge",
        "Where does the existing proxy terminate TLS?",
        "Private endpoints to AHDS",
    ]),
    ("Operations", "Not an afterthought", [
        "APIOps — policies in git, promoted per environment",
        "Named values → contract master lookup",
        "Diagnostics to Log Analytics; App Insights sampling",
        "Who owns policy review and release?",
    ]),
]
for i, (h, sub, items2) in enumerate(cols):
    x = 0.62 + i * 4.06
    card(s, x, 1.6, 3.86, 3.5)
    rect(s, x, 1.6, 3.86, 0.09, fill=TEAL)
    txt(s, x + 0.28, 1.84, 3.3, 0.32, h, size=15, bold=True, color=INK, font=HEAD)
    txt(s, x + 0.28, 2.2, 3.3, 0.26, sub, size=10, color=TEAL, font=MONO, bold=True)
    yy = 2.6
    for it in items2:
        rect(s, x + 0.3, yy + 0.09, 0.07, 0.07, fill=TEAL, shape=MSO_SHAPE.OVAL)
        txt(s, x + 0.5, yy, 3.1, 0.56, it, size=10.5, color=GREY, line_spacing=1.25)
        yy += 0.6

card(s, 0.62, 5.34, 12.1, 1.14, fill=AMBER_LT, line=AMBER)
txt(s, 0.96, 5.54, 11.4, 0.74,
    [("The critical path is provisioning, not engineering.  ", {"bold": True, "color": AMBER}),
     ("Northwind Health already runs APIM elsewhere, so the service is proven internally — but the "
      "request still needs to go in. Estimated on 8/18 at roughly two weeks. Everything in this "
      "deck is deployed and working already; none of it can move to a Northwind Health subscription "
      "until an instance exists.", {"color": INK})], size=11.5, line_spacing=1.3)
notes(s, """
Hand to the APIM specialist / Joe if they are on the call. Otherwise cover it yourself.

SAY:
"Three columns of decisions, and then the one thing that actually gates us.

SKU. The POC runs on BasicV2, which has no SLA and is fine for a proof. For production the
choice is driven almost entirely by network mode, not by throughput. If you need the gateway
injected into a VNet, that's Premium. If public with private endpoints to the back end is
acceptable, StandardV2 is considerably cheaper.

Network. This is really a conversation with your proxy and network team, not with me. The
question I'd put to them: where does TLS terminate today, and does the APIM gateway sit in front
of or behind that? That answer determines the SKU.

Operations. I'd flag the first and last bullets. APIOps means the policies live in git and get
promoted through environments like any other code — they are code, they just happen to be XML.
And 'who owns policy review' is a real organisational question. These policies are the security
boundary. If anyone with portal access can edit them in place, you don't have a boundary, you
have a suggestion.

Now the amber box, and this is the actual message. None of this is blocked on engineering. It's
built and it's running. It's blocked on getting an APIM instance provisioned in a Northwind Health
subscription. Ashay estimated two weeks on the eighteenth even given that you already use APIM
elsewhere. That's the critical path, and it starts the day someone files the request."
""")

# ------------------------------------------------------------------ 16 decisions
s = new(dark=True)
head(s, "Decisions", "Six to own before we leave", dark=True)

decs = [
    ("Dedicated APIM instance, or onboard to the existing estate?",
     "Shared is cheaper; dedicated keeps PHI routing policy under one owner."),
    ("Which SKU for production?", "Follows from network mode. Premium only if VNet injection is required."),
    ("External or internal VNet mode, and where does the proxy sit?",
     "Needs the network team in the room."),
    ("Client secret, private_key_jwt, or mTLS for payers?",
     "Constrained by what your payers can actually support."),
    ("External ID tenant, or the Northwind Health enterprise tenant?",
     "Strong recommendation: External ID."),
    ("Who owns the policy lifecycle and the release gate?",
     "These policies are the security boundary; portal editing should not be the path."),
]
y = 1.50
for i, (q, note) in enumerate(decs):
    card(s, 0.62, y, 12.1, 0.74, fill=NAVY_CARD, line=SLATE)
    chip(s, 0.9, y + 0.18, 0.38, str(i + 1), fill=TEAL_LT, color=NAVY)
    txt(s, 1.52, y + 0.12, 8.2, 0.30, q, size=12.5, bold=True, color=WHITE)
    txt(s, 1.52, y + 0.42, 9.2, 0.26, note, size=10, color=MUTED)
    rect(s, 10.9, y + 0.20, 1.6, 0.34, fill=RGBColor(0x08, 0x1B, 0x30),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25, line=TEAL)
    txt(s, 10.9, y + 0.27, 1.6, 0.24, "owner?", size=9.5, bold=True, color=TEAL_LT,
        align=PP_ALIGN.CENTER, font=MONO)
    y += 0.82

txt(s, 0.62, 6.46, 12.1, 0.34,
    "Owners today. Answers can follow.", size=14, bold=True, color=TEAL_LT, font=HEAD)
notes(s, """
Do not read all six aloud — put the slide up and work the room.

SAY:
"Six decisions. I'm not expecting answers to all of these in the next ninety seconds. I'm asking
for a name against each one, because these are the things that will stall in three weeks if
nobody owns them.

Let me go quickly and you tell me who."

Then go one by one and WRITE THE NAME ON THE SLIDE or in chat. Say the name back out loud so it
is captured on the recording.

If the room goes quiet on number one — dedicated vs shared — offer the tiebreaker:
"The question I'd ask is whether your existing APIM estate has an owner who wants to be
accountable for PHI routing policy. If not, dedicated is the cleaner answer."

If they push back on number six as premature: "It's the one that ages worst. Every month it goes
unowned, more policy drifts into the portal by hand."
""")

# ------------------------------------------------------------------ 17 actions
s = new()
head(s, "Next", "Actions out of this session")

acts = [
    ("Now", "File the APIM provisioning request", "Northwind Health · platform + network",
     "The critical path. Two weeks even with prior usage.", AMBER),
    ("This week", "Confirm payer credential direction", "Northwind Health · integration",
     "Ask two payers what they support: secret, private_key_jwt, or mTLS.", TEAL),
    ("This week", "Network design session with the proxy team", "Joint",
     "Decides SKU. Thirty minutes with the right people in the room.", TEAL),
    ("On instance", "Redeploy the reference implementation", "Microsoft",
     "Bicep and policies are parameterised. Under a day once a target exists.", GREEN),
    ("Follow-up", "Hands-on lab against the deployed gateway", "Joint",
     "Your engineers run the tests and edit a policy themselves.", GREEN),
]
y = 1.55
for when, what, who, why, col in acts:
    card(s, 0.62, y, 12.1, 0.82)
    rect(s, 0.62, y, 0.075, 0.82, fill=col)
    rect(s, 0.92, y + 0.21, 1.5, 0.4, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.25, line=col)
    txt(s, 0.92, y + 0.30, 1.5, 0.24, when, size=9.5, bold=True, color=col,
        align=PP_ALIGN.CENTER, font=MONO)
    txt(s, 2.68, y + 0.13, 6.6, 0.3, what, size=13.5, bold=True, color=INK)
    txt(s, 2.68, y + 0.45, 6.8, 0.26, why, size=10.5, color=GREY)
    txt(s, 9.6, y + 0.26, 2.9, 0.3, who, size=10, bold=True, color=TEAL,
        align=PP_ALIGN.RIGHT, font=MONO)
    y += 0.88

card(s, 0.62, 6.02, 12.1, 0.72, fill=RGBColor(0x0B, 0x22, 0x3A), line=None)
txt(s, 0.96, 6.24, 11.5, 0.32,
    [("Everything shown today is deployed and reproducible: Bicep, policies, and a test suite "
      "whose assertions are the security guarantees.", {"color": ICE, "bold": True})], size=12)
notes(s, """
SAY:
"Five actions, and only the first one is on the critical path.

File the provisioning request. That's the two-week clock and nothing else starts until it does.
Ashay, on the eighteenth you mentioned working with Sunil and Surya and the proxy team — is that
still the right route?

Second, and you can do this in parallel: go and ask two of your payers what they can actually
support for credentials. That answer constrains our design more than our preference does, and
it's a phone call.

Third, thirty minutes with your network team to settle the SKU question.

Fourth is ours. Once an instance exists in a Northwind Health subscription, redeploying this is under
a day — the Bicep and the policies are parameterised, there's nothing bespoke about the
environment you just watched.

Fifth is the one I'd actually push for. Rather than another briefing, let's do a hands-on lab
where your engineers run the test suite themselves and edit a policy. You'll learn more in
ninety minutes of that than in three of these.

Closing line: everything you saw today is deployed and reproducible. It's Bicep you can read,
policies you can attach to a debug trace, and a test suite where the assertions *are* the
security guarantees. If those stop being true, the suite goes red. That's the standard I'd want
you to hold this to."
""")

out = Path(__file__).with_name("Prov-Azure-APIM-Briefing.pptx")
prs.core_properties.title = "Northwind Health — Azure API Management Briefing"
prs.core_properties.author = "Microsoft Health & Life Sciences"
prs.save(out)
print(f"{len(prs.slides._sldIdLst)} slides -> {out}")
